"""
Cantio - Presentation Editor v2
QGraphicsScene/View slide editor with: Text, Image, Shape (rect/ellipse/line/
triangle/star/arrow), Code, Table elements; rotation, shadow, gradient fills,
gradient backgrounds, grid/snap, smart guides, animation timeline, slide notes,
templates, PPTX import, undo/redo.
"""
from __future__ import annotations

import os, json, copy, math

from PyQt6.QtWidgets import (
    QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, QSplitter,
    QPushButton, QLabel, QListWidget, QListWidgetItem, QScrollArea,
    QFormLayout, QLineEdit, QSpinBox, QDoubleSpinBox, QColorDialog,
    QFileDialog, QMessageBox, QComboBox, QCheckBox, QFrame,
    QToolBar, QSizePolicy, QInputDialog, QDialog, QDialogButtonBox,
    QTabWidget, QTextEdit, QGraphicsScene, QGraphicsView, QGraphicsItem,
    QGraphicsRectItem, QGraphicsEllipseItem, QApplication, QSlider,
    QGroupBox, QScrollBar, QStackedWidget, QToolButton, QFontComboBox,
    QGraphicsDropShadowEffect, QGridLayout,
)
from PyQt6.QtCore import (
    Qt, QSize, QRect, QPoint, QPointF, QRectF, QSizeF,
    pyqtSignal, QTimer,
)
from PyQt6.QtGui import (
    QFont, QColor, QPainter, QPixmap, QPen, QBrush,
    QLinearGradient, QRadialGradient, QAction, QKeySequence,
    QCursor, QTransform, QPainterPath, QPolygonF,
    QSyntaxHighlighter, QTextCharFormat,
)

# ── Constants ─────────────────────────────────────────────────────────────────

CANVAS_W, CANVAS_H = 1920, 1080
HANDLE_SIZE = 10
MAX_UNDO    = 50

ET_TEXT     = "text"
ET_IMAGE    = "image"
ET_RECT     = "rect"
ET_ELLIPSE  = "ellipse"
ET_LINE     = "line"
ET_TRIANGLE = "triangle"
ET_STAR     = "star"
ET_ARROW    = "arrow"
ET_CODE     = "code"
ET_TABLE    = "table"
ET_CHART    = "chart"
ET_DIAGRAM  = "diagram"

_SHAPE_TYPES = (ET_RECT, ET_ELLIPSE, ET_LINE, ET_TRIANGLE, ET_STAR, ET_ARROW)

_ENTRANCES = ["none","fade_in","slide_left","slide_right",
              "slide_up","slide_down","zoom_in","bounce","blur_in"]
_EXITS     = ["none","fade_out","slide_left","slide_right",
              "slide_up","slide_down","zoom_out","blur_out"]
_LOOPS     = ["none","pulse","float","glow","shake"]
_TRANSITIONS = ["none","fade","slide_left","slide_right",
                "zoom_in","blur","instant"]

_STYLE = """
QMainWindow,QWidget{background:#181818;color:#e0e0e0;
    font-family:'Segoe UI',sans-serif;font-size:12px;}
QPushButton{background:#232323;color:#e0e0e0;border:1px solid #2c2c2c;
    border-radius:4px;padding:5px 10px;}
QPushButton:hover{background:#2a2a2a;border-color:#3a3a3a;}
QPushButton:checked{background:#1a3a5c;border-color:#5294e2;color:#5294e2;}
QPushButton:disabled{color:#444;border-color:#222;}
QLineEdit,QSpinBox,QDoubleSpinBox,QComboBox,QTextEdit,QFontComboBox{
    background:#1c1c1c;color:#e0e0e0;border:1px solid #262626;
    border-radius:4px;padding:4px 6px;}
QLineEdit:focus,QSpinBox:focus{border-color:#5294e2;}
QListWidget{background:#141414;color:#e0e0e0;border:none;}
QListWidget::item{padding:6px 8px;border-radius:3px;margin:1px 2px;}
QListWidget::item:hover{background:#1e1e1e;}
QListWidget::item:selected{background:#1c3a5a;}
QLabel{color:#cccccc;}
QToolBar{background:#0f0f0f;border-bottom:1px solid #1c1c1c;
    padding:2px 4px;spacing:3px;}
QGraphicsView{border:none;background:#111;}
QScrollArea{border:none;}
QGroupBox{border:1px solid #2a2a2a;border-radius:4px;
    margin-top:8px;padding-top:4px;font-size:11px;color:#888;}
QGroupBox::title{subcontrol-origin:margin;left:8px;}
QSplitter::handle{background:#1c1c1c;}
QCheckBox{color:#ccc;}
QSlider::groove:horizontal{background:#2a2a2a;height:4px;border-radius:2px;}
QSlider::handle:horizontal{background:#5294e2;width:14px;height:14px;
    border-radius:7px;margin:-5px 0;}
"""

# ── Element default factories ─────────────────────────────────────────────────

def _default_element(kind: str) -> dict:
    base = {
        "type": kind, "x": 200, "y": 200, "w": 400, "h": 100,
        "z": 0, "locked": False, "visible": True, "rotation": 0,
        "shadow": {"enabled": False, "color": "#000000",
                   "blur": 8, "offset_x": 3, "offset_y": 3},
        "animation": {"entrance": "none", "exit": "none",
                      "loop": "none", "delay": 0, "duration": 500},
    }
    if kind == ET_TEXT:
        base.update({
            "text": "Text nou", "font": "Segoe UI", "font_size": 48,
            "bold": False, "italic": False, "underline": False,
            "color": "#ffffff", "align": "center",
            "line_spacing": 1.2, "letter_spacing": 0,
            "bg_color": "", "border_color": "",
            "h": 120,
        })
    elif kind == ET_IMAGE:
        base.update({"path": "", "opacity": 1.0, "border_radius": 0,
                     "brightness": 1.0, "contrast": 1.0})
    elif kind in _SHAPE_TYPES:
        base.update({
            "fill": "#5294e2", "border_color": "#ffffff",
            "border_width": 2, "opacity": 1.0, "border_radius": 0,
            "fill_type": "solid",
            "gradient_from": "#5294e2", "gradient_to": "#1a1a5a",
            "gradient_angle": 90,
        })
        if kind == ET_ELLIPSE:
            base["fill"] = "#a6e3a1"
        elif kind == ET_LINE:
            base.update({"color": "#ffffff", "line_width": 3, "h": 3})
        elif kind == ET_TRIANGLE:
            base["fill"] = "#f38ba8"
        elif kind == ET_STAR:
            base.update({"fill": "#f9e2af", "points": 5})
        elif kind == ET_ARROW:
            base.update({"fill": "#89dceb", "direction": "right"})
    elif kind == ET_CODE:
        base.update({
            "code_text": "# Cod sursă\nprint('Hello, World!')",
            "language": "python",
            "font_size": 20, "bg_color": "#1e1e2e",
            "text_color": "#cdd6f4", "w": 700, "h": 300,
        })
    elif kind == ET_TABLE:
        base.update({
            "rows": 3, "cols": 3,
            "cells": [["Cap 1","Cap 2","Cap 3"],
                      ["R1C1","R1C2","R1C3"],
                      ["R2C1","R2C2","R2C3"]],
            "header_bg": "#1a3a5a", "header_color": "#ffffff",
            "cell_bg": "#1c1c1c", "cell_color": "#e0e0e0",
            "border_color": "#333333", "font_size": 20,
            "w": 700, "h": 250,
        })
    elif kind == ET_CHART:
        base.update({
            "chart_type": "bar",
            "title": "Grafic",
            "labels": ["Ian", "Feb", "Mar", "Apr", "Mai"],
            "values": [40.0, 65.0, 50.0, 80.0, 55.0],
            "value_color": "#5294e2",
            "bg_color": "#1c1c1c",
            "text_color": "#e0e0e0",
            "grid_color": "#2a2a2a",
            "show_values": True,
            "w": 700, "h": 400,
        })
    elif kind == ET_DIAGRAM:
        base.update({
            "nodes": ["Root", "Ramura A", "Ramura B", "Sub A1", "Sub B1"],
            "edges": [[0, 1], [0, 2], [1, 3], [2, 4]],
            "root_color": "#5294e2",
            "node_color": "#1a3a5c",
            "line_color": "#888888",
            "text_color": "#ffffff",
            "font_size": 18,
            "w": 700, "h": 450,
        })
    return base


def _default_slide() -> dict:
    return {
        "bg_type": "solid",
        "bg_color": "#000000",
        "bg_image": "",
        "bg_gradient_from": "#1a1a2a",
        "bg_gradient_to":   "#000000",
        "bg_gradient_angle": 135,
        "elements": [],
        "transition": "fade",
        "transition_ms": 400,
        "notes": "",
    }


# ── Syntax highlighter for code elements ──────────────────────────────────────

class _CodeHighlighter(QSyntaxHighlighter):
    _RULES = {
        "python": {
            "keyword": (r'\b(def|class|import|from|return|if|elif|else|for|'
                        r'while|try|except|finally|with|as|pass|break|'
                        r'continue|lambda|yield|and|or|not|in|is|True|False|None)\b',
                        "#cba6f7"),
            "builtin": (r'\b(print|len|range|int|str|float|list|dict|set|'
                        r'tuple|type|open|input|super|self)\b', "#89b4fa"),
            "string":  (r'(\".*?\"|\'.*?\'|\"\"\"[\s\S]*?\"\"\"|'
                        r"\'\'\'[\s\S]*?\'\'\')", "#a6e3a1"),
            "comment": (r'#[^\n]*', "#6c7086"),
            "number":  (r'\b\d+(\.\d+)?\b', "#fab387"),
        },
        "javascript": {
            "keyword": (r'\b(var|let|const|function|return|if|else|for|while|'
                        r'try|catch|finally|class|import|export|from|async|'
                        r'await|new|this|typeof|instanceof|true|false|null)\b',
                        "#cba6f7"),
            "string":  (r'(\".*?\"|\'.*?\'|`.*?`)', "#a6e3a1"),
            "comment": (r'//[^\n]*', "#6c7086"),
            "number":  (r'\b\d+(\.\d+)?\b', "#fab387"),
        },
    }

    def __init__(self, document, language="python"):
        super().__init__(document)
        import re
        self._formats = []
        rules = self._RULES.get(language, self._RULES["python"])
        for _, (pattern, color) in rules.items():
            fmt = QTextCharFormat()
            fmt.setForeground(QColor(color))
            self._formats.append((re.compile(pattern), fmt))

    def highlightBlock(self, text: str):
        for regex, fmt in self._formats:
            for m in regex.finditer(text):
                self.setFormat(m.start(), m.end() - m.start(), fmt)


# ══════════════════════════════════════════════════════════════════════════════
# Handles
# ══════════════════════════════════════════════════════════════════════════════

class ResizeHandle(QGraphicsRectItem):
    CURSORS = {
        "tl": Qt.CursorShape.SizeFDiagCursor, "tr": Qt.CursorShape.SizeBDiagCursor,
        "bl": Qt.CursorShape.SizeBDiagCursor, "br": Qt.CursorShape.SizeFDiagCursor,
        "t":  Qt.CursorShape.SizeVerCursor,  "b":  Qt.CursorShape.SizeVerCursor,
        "l":  Qt.CursorShape.SizeHorCursor,  "r":  Qt.CursorShape.SizeHorCursor,
    }

    def __init__(self, pos_key: str, parent):
        s = HANDLE_SIZE
        super().__init__(-s//2, -s//2, s, s, parent)
        self.pos_key = pos_key
        self.setFlag(QGraphicsItem.GraphicsItemFlag.ItemIsMovable, False)
        self.setBrush(QBrush(QColor("#5294e2")))
        self.setPen(QPen(QColor("#ffffff"), 1))
        self.setCursor(QCursor(self.CURSORS.get(pos_key, Qt.CursorShape.SizeAllCursor)))
        self.setZValue(9999)

    def mousePressEvent(self, e): e.accept()

    def mouseMoveEvent(self, e):
        parent = self.parentItem()
        if parent and hasattr(parent, '_resize_by_handle'):
            parent._resize_by_handle(self.pos_key, e.pos() - e.lastPos())

    def mouseReleaseEvent(self, e):
        parent = self.parentItem()
        if parent and hasattr(parent, '_on_resize_done'):
            parent._on_resize_done()
        e.accept()


class RotationHandle(QGraphicsEllipseItem):
    _R = 7

    def __init__(self, parent):
        r = self._R
        super().__init__(-r, -r, r*2, r*2, parent)
        self._parent = parent
        self._dragging = False
        self.setBrush(QBrush(QColor("#ffffff")))
        self.setPen(QPen(QColor("#5294e2"), 2))
        self.setCursor(QCursor(Qt.CursorShape.CrossCursor))
        self.setZValue(10000)
        self._reposition()

    def _reposition(self):
        self.setPos(self._parent.data["w"] / 2, -28)

    def mousePressEvent(self, e):
        self._dragging = True
        e.accept()

    def mouseMoveEvent(self, e):
        if not self._dragging:
            return
        cx = self._parent.data["w"] / 2
        cy = self._parent.data["h"] / 2
        mp = self._parent.mapFromScene(
            self._parent.mapToScene(e.pos() + self.pos()))
        dx = mp.x() - cx
        dy = mp.y() - cy
        angle = (math.degrees(math.atan2(dy, dx)) + 90) % 360
        self._parent.data["rotation"] = round(angle, 1)
        self._parent.setRotation(angle)
        if self._parent.scene_ref:
            self._parent.scene_ref._on_element_changed()
        e.accept()

    def mouseReleaseEvent(self, e):
        self._dragging = False
        e.accept()


# ══════════════════════════════════════════════════════════════════════════════
# Base Element
# ══════════════════════════════════════════════════════════════════════════════

class BaseElement(QGraphicsItem):
    def __init__(self, data: dict, scene_ref):
        super().__init__()
        self.data = data
        self.scene_ref = scene_ref
        self._handles: dict[str, ResizeHandle] = {}
        self._rot_handle: RotationHandle | None = None
        self._selected_locally = False

        self.setFlag(QGraphicsItem.GraphicsItemFlag.ItemIsMovable)
        self.setFlag(QGraphicsItem.GraphicsItemFlag.ItemIsSelectable)
        self.setFlag(QGraphicsItem.GraphicsItemFlag.ItemSendsGeometryChanges)
        self.setFlag(QGraphicsItem.GraphicsItemFlag.ItemIsFocusable)
        self.setPos(data["x"], data["y"])
        self.setZValue(data.get("z", 0))
        self.setVisible(data.get("visible", True))
        self.setOpacity(float(data.get("opacity", 1.0)))
        rot = data.get("rotation", 0)
        if rot:
            self.setTransformOriginPoint(data["w"] / 2, data["h"] / 2)
            self.setRotation(rot)
        self._build_handles()
        self._apply_shadow()

    def _build_handles(self):
        for key in ("tl","t","tr","l","r","bl","b","br"):
            h = ResizeHandle(key, self)
            h.hide()
            self._handles[key] = h
        self._rot_handle = RotationHandle(self)
        self._rot_handle.hide()
        self._update_handle_positions()

    def _update_handle_positions(self):
        w, h = self.data["w"], self.data["h"]
        self.setTransformOriginPoint(w / 2, h / 2)
        positions = {
            "tl": QPointF(0,   0),   "t":  QPointF(w/2, 0),
            "tr": QPointF(w,   0),   "l":  QPointF(0,   h/2),
            "r":  QPointF(w,   h/2), "bl": QPointF(0,   h),
            "b":  QPointF(w/2, h),   "br": QPointF(w,   h),
        }
        for key, pos in positions.items():
            if key in self._handles:
                self._handles[key].setPos(pos)
        if self._rot_handle:
            self._rot_handle._reposition()

    def _apply_shadow(self):
        sh = self.data.get("shadow", {})
        if sh.get("enabled"):
            fx = QGraphicsDropShadowEffect()
            fx.setBlurRadius(float(sh.get("blur", 8)))
            fx.setColor(QColor(sh.get("color", "#000000")))
            fx.setOffset(float(sh.get("offset_x", 3)),
                         float(sh.get("offset_y", 3)))
            self.setGraphicsEffect(fx)
        else:
            self.setGraphicsEffect(None)

    def set_selected_state(self, sel: bool):
        self._selected_locally = sel
        locked = self.data.get("locked", False)
        for h in self._handles.values():
            h.setVisible(sel and not locked)
        if self._rot_handle:
            self._rot_handle.setVisible(sel and not locked)
        self.update()

    def boundingRect(self):
        return QRectF(0, 0, self.data["w"], self.data["h"])

    def itemChange(self, change, value):
        if change == QGraphicsItem.GraphicsItemChange.ItemPositionChange:
            if self.scene_ref and getattr(self.scene_ref, '_snap_enabled', False):
                gs = getattr(self.scene_ref, '_grid_size', 20)
                nx = round(value.x() / gs) * gs
                ny = round(value.y() / gs) * gs
                value = QPointF(nx, ny)
        if change == QGraphicsItem.GraphicsItemChange.ItemPositionHasChanged:
            self.data["x"] = round(self.pos().x())
            self.data["y"] = round(self.pos().y())
            if self.scene_ref and hasattr(self.scene_ref, '_on_element_changed'):
                self.scene_ref._on_element_changed()
        return super().itemChange(change, value)

    def _resize_by_handle(self, key: str, delta: QPointF):
        dx, dy = delta.x(), delta.y()
        x, y, w, h = (self.data["x"], self.data["y"],
                      self.data["w"], self.data["h"])
        if "l" in key: x += dx; w -= dx
        if "r" in key: w += dx
        if "t" in key: y += dy; h -= dy
        if "b" in key: h += dy
        w, h = max(20, w), max(10, h)
        self.data.update({"x": int(x), "y": int(y),
                          "w": int(w), "h": int(h)})
        self.setPos(x, y)
        self.prepareGeometryChange()
        self._update_handle_positions()
        self.update()

    def _on_resize_done(self):
        if self.scene_ref:
            self.scene_ref._on_element_changed()

    def keyPressEvent(self, e):
        step = 10 if not (e.modifiers() & Qt.KeyboardModifier.ShiftModifier) else 1
        k = e.key()
        dx = dy = 0
        if k == Qt.Key.Key_Left:  dx = -step
        elif k == Qt.Key.Key_Right: dx = step
        elif k == Qt.Key.Key_Up:   dy = -step
        elif k == Qt.Key.Key_Down: dy = step
        if dx or dy:
            self.moveBy(dx, dy)
        else:
            super().keyPressEvent(e)

    def contextMenuEvent(self, e):
        from PyQt6.QtWidgets import QMenu
        menu = QMenu()
        menu.addAction("🗑 Șterge",         lambda: self._delete_self())
        menu.addAction("🔝 Aduce în față",  lambda: self._bring_forward())
        menu.addAction("🔙 Trimite în spate", lambda: self._send_back())
        locked = self.data.get("locked", False)
        menu.addAction("🔓 Deblochează" if locked else "🔒 Blochează",
                       lambda: self._toggle_lock())
        menu.addSeparator()
        menu.addAction("📋 Duplică", lambda: self._duplicate_self())
        menu.exec(e.screenPos())

    def _delete_self(self):
        if self.scene():
            self.scene().removeItem(self)
            if self.scene_ref:
                self.scene_ref._on_element_deleted(self.data)

    def _duplicate_self(self):
        if self.scene_ref:
            new_data = copy.deepcopy(self.data)
            new_data["x"] += 20; new_data["y"] += 20
            self.scene_ref._slide_data.setdefault("elements", []).append(new_data)
            item = _make_element(new_data, self.scene_ref)
            self.scene_ref.addItem(item)
            self.scene_ref._items.append(item)
            self.scene_ref._select_item(item)

    def _bring_forward(self):
        self.data["z"] = self.data.get("z", 0) + 1
        self.setZValue(self.data["z"])

    def _send_back(self):
        self.data["z"] = self.data.get("z", 0) - 1
        self.setZValue(self.data["z"])

    def _toggle_lock(self):
        self.data["locked"] = not self.data.get("locked", False)
        locked = self.data["locked"]
        self.setFlag(QGraphicsItem.GraphicsItemFlag.ItemIsMovable, not locked)
        for h in self._handles.values():
            h.setVisible(self._selected_locally and not locked)


# ══════════════════════════════════════════════════════════════════════════════
# Concrete element types
# ══════════════════════════════════════════════════════════════════════════════

class TextElement(BaseElement):
    def paint(self, p, option, widget=None):
        d = self.data
        w, h = d["w"], d["h"]
        bg = d.get("bg_color", "")
        if bg:
            p.fillRect(QRectF(0, 0, w, h), QColor(bg))
        bc = d.get("border_color", "")
        if bc:
            p.setPen(QPen(QColor(bc), 2))
            p.setBrush(Qt.BrushStyle.NoBrush)
            p.drawRect(QRectF(1, 1, w-2, h-2))
        font = QFont(d.get("font", "Segoe UI"), int(d.get("font_size", 48)))
        font.setBold(d.get("bold", False))
        font.setItalic(d.get("italic", False))
        font.setUnderline(d.get("underline", False))
        ls = d.get("letter_spacing", 0)
        if ls:
            font.setLetterSpacing(QFont.SpacingType.AbsoluteSpacing, ls)
        p.setFont(font)
        p.setPen(QColor(d.get("color", "#ffffff")))
        align_map = {
            "left":   Qt.AlignmentFlag.AlignLeft,
            "center": Qt.AlignmentFlag.AlignHCenter,
            "right":  Qt.AlignmentFlag.AlignRight,
        }
        align = align_map.get(d.get("align", "center"),
                               Qt.AlignmentFlag.AlignHCenter)
        flags = align | Qt.AlignmentFlag.AlignVCenter | Qt.TextFlag.TextWordWrap
        p.drawText(QRectF(4, 4, w-8, h-8), flags, d.get("text", ""))
        if self._selected_locally:
            p.setPen(QPen(QColor("#5294e2"), 2, Qt.PenStyle.DashLine))
            p.setBrush(Qt.BrushStyle.NoBrush)
            p.drawRect(QRectF(0, 0, w, h))

    def mouseDoubleClickEvent(self, e):
        text, ok = QInputDialog.getMultiLineText(
            None, "Editare text", "Text:", self.data.get("text", ""))
        if ok:
            self.data["text"] = text
            self.update()
            if self.scene_ref:
                self.scene_ref._on_element_changed()


class ImageElement(BaseElement):
    def __init__(self, data, scene_ref):
        super().__init__(data, scene_ref)
        self._pixmap: QPixmap | None = None
        self._load_pixmap()

    def _load_pixmap(self):
        path = self.data.get("path", "")
        self._pixmap = QPixmap(path) if path and os.path.exists(path) else None

    def paint(self, p, option, widget=None):
        d = self.data
        w, h = d["w"], d["h"]
        p.setOpacity(float(d.get("opacity", 1.0)))
        if self._pixmap and not self._pixmap.isNull():
            scaled = self._pixmap.scaled(
                int(w), int(h),
                Qt.AspectRatioMode.KeepAspectRatio,
                Qt.TransformationMode.SmoothTransformation)
            ox = (w - scaled.width()) / 2
            oy = (h - scaled.height()) / 2
            p.drawPixmap(int(ox), int(oy), scaled)
        else:
            p.fillRect(QRectF(0, 0, w, h), QColor("#333333"))
            p.setPen(QColor("#888888"))
            p.drawText(QRectF(0, 0, w, h), Qt.AlignmentFlag.AlignCenter,
                       "🖼 Imagine\n(fără cale)")
        p.setOpacity(1.0)
        if self._selected_locally:
            p.setPen(QPen(QColor("#5294e2"), 2, Qt.PenStyle.DashLine))
            p.setBrush(Qt.BrushStyle.NoBrush)
            p.drawRect(QRectF(0, 0, w, h))


def _shape_gradient(d: dict, w: float, h: float) -> QBrush:
    angle = d.get("gradient_angle", 90)
    rad = math.radians(angle)
    cx, cy = w / 2, h / 2
    dx, dy = math.cos(rad) * w / 2, math.sin(rad) * h / 2
    grad = QLinearGradient(cx - dx, cy - dy, cx + dx, cy + dy)
    grad.setColorAt(0, QColor(d.get("gradient_from", "#5294e2")))
    grad.setColorAt(1, QColor(d.get("gradient_to", "#1a1a5a")))
    return QBrush(grad)


def _build_star_path(cx, cy, outer, inner, points=5) -> QPainterPath:
    path = QPainterPath()
    for i in range(points * 2):
        r = outer if i % 2 == 0 else inner
        angle = math.pi / points * i - math.pi / 2
        x = cx + r * math.cos(angle)
        y = cy + r * math.sin(angle)
        if i == 0:
            path.moveTo(x, y)
        else:
            path.lineTo(x, y)
    path.closeSubpath()
    return path


def _build_arrow_path(w, h, direction="right") -> QPainterPath:
    path = QPainterPath()
    hw = h * 0.35
    nw = w * 0.60
    if direction == "right":
        path.moveTo(0, (h - hw) / 2)
        path.lineTo(nw, (h - hw) / 2)
        path.lineTo(nw, 0)
        path.lineTo(w, h / 2)
        path.lineTo(nw, h)
        path.lineTo(nw, (h + hw) / 2)
        path.lineTo(0, (h + hw) / 2)
    elif direction == "left":
        path.moveTo(w, (h - hw) / 2)
        path.lineTo(w - nw, (h - hw) / 2)
        path.lineTo(w - nw, 0)
        path.lineTo(0, h / 2)
        path.lineTo(w - nw, h)
        path.lineTo(w - nw, (h + hw) / 2)
        path.lineTo(w, (h + hw) / 2)
    path.closeSubpath()
    return path


# ── Chart / Diagram shared drawing helpers ────────────────────────────────────

def _diagram_layout(nodes, edges, w, h):
    """BFS tree layout → list of (x, y) centres for each node."""
    n = len(nodes)
    if n == 0:
        return []
    children = {i: [] for i in range(n)}
    for a, b in edges:
        if 0 <= a < n and 0 <= b < n:
            children[a].append(b)
    visited = [False] * n
    levels: dict[int, int] = {}
    queue = [0]
    visited[0] = True
    lvl = 0
    while queue:
        next_q = []
        for node in queue:
            levels[node] = lvl
            for c in children[node]:
                if not visited[c]:
                    visited[c] = True
                    next_q.append(c)
        queue = next_q
        lvl += 1
    for i in range(n):
        if not visited[i]:
            levels[i] = lvl; lvl += 1
    level_nodes: dict[int, list] = {}
    for node, lv in levels.items():
        level_nodes.setdefault(lv, []).append(node)
    max_lv = max(levels.values()) if levels else 0
    pad = 40.0
    avail_h = h - pad * 2
    lv_h = avail_h / (max_lv + 1) if max_lv > 0 else avail_h
    positions = [None] * n
    for lv, lst in level_nodes.items():
        avail_w = w - pad * 2
        spacing = avail_w / (len(lst) + 1)
        y = pad + lv * lv_h + lv_h / 2
        for j, node in enumerate(sorted(lst)):
            positions[node] = (pad + spacing * (j + 1), y)
    return positions


def _draw_chart(p: QPainter, d: dict, ew: float, eh: float):
    """Render chart data dict onto painter at (0,0) size ew×eh."""
    chart_type = d.get("chart_type", "bar")
    bg = QColor(d.get("bg_color", "#1c1c1c"))
    p.fillRect(QRectF(0, 0, ew, eh), bg)

    text_col = QColor(d.get("text_color", "#e0e0e0"))
    val_col  = QColor(d.get("value_color", "#5294e2"))
    grid_col = QColor(d.get("grid_color", "#2a2a2a"))
    show_vals = d.get("show_values", True)
    title    = d.get("title", "")
    labels   = list(d.get("labels", []))
    values   = [float(v) for v in d.get("values", []) if str(v).replace(".","").replace("-","").isdigit() or isinstance(v, (int,float))]

    title_h = 0.0
    if title:
        font = QFont("Segoe UI", max(6, int(eh * 0.055)))
        font.setBold(True)
        p.setFont(font)
        p.setPen(text_col)
        title_h = eh * 0.12
        p.drawText(QRectF(0, 4, ew, title_h),
                   Qt.AlignmentFlag.AlignHCenter | Qt.AlignmentFlag.AlignVCenter, title)

    if not values:
        p.setPen(text_col)
        p.drawText(QRectF(0, 0, ew, eh), Qt.AlignmentFlag.AlignCenter, "Fără date")
        return

    n = len(values)
    labels = (labels + [""] * n)[:n]

    if chart_type == "pie":
        _draw_chart_pie(p, ew, eh, title_h, labels, values, val_col, text_col, show_vals)
    elif chart_type == "line":
        _draw_chart_line(p, ew, eh, title_h, labels, values, val_col, text_col, grid_col, show_vals)
    else:
        _draw_chart_bar(p, ew, eh, title_h, labels, values, val_col, text_col, grid_col, show_vals)


def _draw_chart_bar(p, ew, eh, title_h, labels, values, val_col, text_col, grid_col, show_vals):
    n = len(values)
    max_val = max(values) if values else 1
    pad_l = ew * 0.09; pad_r = ew * 0.03
    pad_t = title_h + eh * 0.05; pad_b = eh * 0.18
    cw = ew - pad_l - pad_r; ch = eh - pad_t - pad_b
    # Grid & y-labels
    font_sm = QFont("Segoe UI", max(5, int(eh * 0.038)))
    p.setFont(font_sm)
    for i in range(5):
        gy = pad_t + ch * (1 - i / 4)
        p.setPen(QPen(grid_col, 1)); p.drawLine(QPointF(pad_l, gy), QPointF(pad_l + cw, gy))
        p.setPen(text_col)
        p.drawText(QRectF(0, gy - eh*0.03, pad_l - 2, eh*0.06),
                   Qt.AlignmentFlag.AlignRight | Qt.AlignmentFlag.AlignVCenter,
                   f"{max_val*i/4:.0f}")
    # Axes
    p.setPen(QPen(QColor("#555555"), max(1, int(eh*0.005))))
    p.drawLine(QPointF(pad_l, pad_t), QPointF(pad_l, pad_t+ch))
    p.drawLine(QPointF(pad_l, pad_t+ch), QPointF(pad_l+cw, pad_t+ch))
    # Bars
    slot = cw / n; bw = slot * 0.6
    for i, (lbl, val) in enumerate(zip(labels, values)):
        bh = (val / max_val) * ch if max_val else 0
        bx = pad_l + slot * i + slot * 0.2; by = pad_t + ch - bh
        if bh > 0:
            grad = QLinearGradient(bx, by, bx, by + bh)
            grad.setColorAt(0, val_col.lighter(130)); grad.setColorAt(1, val_col)
            p.fillRect(QRectF(bx, by, bw, bh), QBrush(grad))
        p.setFont(font_sm); p.setPen(text_col)
        p.drawText(QRectF(bx - slot*0.1, pad_t+ch+3, bw+slot*0.2, pad_b-3),
                   Qt.AlignmentFlag.AlignHCenter | Qt.AlignmentFlag.AlignTop, str(lbl))
        if show_vals and bh > 0:
            vstr = f"{val:.0f}" if val == int(val) else f"{val:.1f}"
            p.drawText(QRectF(bx, max(by - eh*0.07, 0), bw, eh*0.07),
                       Qt.AlignmentFlag.AlignHCenter | Qt.AlignmentFlag.AlignVCenter, vstr)


def _draw_chart_line(p, ew, eh, title_h, labels, values, val_col, text_col, grid_col, show_vals):
    n = len(values)
    max_val = max(values) if values else 1
    pad_l = ew * 0.09; pad_r = ew * 0.03
    pad_t = title_h + eh * 0.05; pad_b = eh * 0.18
    cw = ew - pad_l - pad_r; ch = eh - pad_t - pad_b
    font_sm = QFont("Segoe UI", max(5, int(eh * 0.038)))
    # Grid
    for i in range(5):
        gy = pad_t + ch * (1 - i / 4)
        p.setPen(QPen(grid_col, 1)); p.drawLine(QPointF(pad_l, gy), QPointF(pad_l+cw, gy))
        p.setFont(font_sm); p.setPen(text_col)
        p.drawText(QRectF(0, gy-eh*0.03, pad_l-2, eh*0.06),
                   Qt.AlignmentFlag.AlignRight | Qt.AlignmentFlag.AlignVCenter, f"{max_val*i/4:.0f}")
    p.setPen(QPen(QColor("#555555"), max(1, int(eh*0.005))))
    p.drawLine(QPointF(pad_l, pad_t), QPointF(pad_l, pad_t+ch))
    p.drawLine(QPointF(pad_l, pad_t+ch), QPointF(pad_l+cw, pad_t+ch))
    spacing = cw / (n - 1) if n > 1 else cw
    pts = [QPointF(pad_l + spacing*i, pad_t + ch*(1 - v/max_val)) for i, v in enumerate(values)]
    # Line
    p.setPen(QPen(val_col, max(2, int(eh*0.012))))
    for i in range(len(pts)-1):
        p.drawLine(pts[i], pts[i+1])
    # Points
    r2 = max(3.0, eh * 0.018)
    p.setBrush(QBrush(val_col)); p.setPen(QPen(QColor("#ffffff"), max(1, int(r2*0.4))))
    for pt in pts:
        p.drawEllipse(pt, r2, r2)
    # Labels
    slot = cw / n if n else cw
    p.setFont(font_sm); p.setPen(text_col)
    for i, lbl in enumerate(labels):
        px = pad_l + spacing*i if n > 1 else pad_l + cw/2
        p.drawText(QRectF(px - slot/2, pad_t+ch+3, slot, pad_b-3),
                   Qt.AlignmentFlag.AlignHCenter | Qt.AlignmentFlag.AlignTop, str(lbl))


def _draw_chart_pie(p, ew, eh, title_h, labels, values, val_col, text_col, show_vals):
    total = sum(values) if values else 1
    pie_area_h = eh - title_h
    pie_size = min(ew * 0.55, pie_area_h * 0.85)
    cx = ew * 0.33; cy = title_h + pie_area_h / 2
    r = pie_size / 2
    # Color palette
    base_hue = val_col.hsvHueF()
    colors = [QColor.fromHsvF((base_hue + i * 0.15) % 1.0, 0.7,
                               0.75 + 0.1*(i%2)) for i in range(len(values))]
    start = 90 * 16
    for i, val in enumerate(values):
        span = -int(360 * 16 * val / total)
        p.setBrush(QBrush(colors[i]))
        p.setPen(QPen(QColor("#1c1c1c"), max(1, int(r*0.02))))
        p.drawPie(QRectF(cx - r, cy - r, r*2, r*2), start, span)
        start += span
    # Legend
    font_sm = QFont("Segoe UI", max(5, int(eh * 0.042)))
    p.setFont(font_sm)
    lx = cx + r + 14
    entry_h = min(eh * 0.085, (eh - title_h - 8) / max(len(values), 1))
    for i, (lbl, val) in enumerate(zip(labels, values)):
        ly = title_h + 8 + i * entry_h
        p.fillRect(QRectF(lx, ly + entry_h*0.2, 12, entry_h*0.6), colors[i])
        p.setPen(text_col)
        pct = f" ({val/total*100:.0f}%)" if show_vals else ""
        p.drawText(QRectF(lx+16, ly, ew-lx-16-4, entry_h),
                   Qt.AlignmentFlag.AlignVCenter, f"{lbl}{pct}")


def _draw_diagram(p: QPainter, d: dict, ew: float, eh: float):
    """Render diagram/organigram onto painter at (0,0) size ew×eh."""
    nodes = d.get("nodes", [])
    edges = d.get("edges", [])
    if not nodes:
        p.setPen(QColor(d.get("text_color", "#e0e0e0")))
        p.drawText(QRectF(0, 0, ew, eh), Qt.AlignmentFlag.AlignCenter, "Fără noduri")
        return

    root_col = QColor(d.get("root_color", "#5294e2"))
    node_col = QColor(d.get("node_color", "#1a3a5c"))
    line_col = QColor(d.get("line_color", "#888888"))
    text_col = QColor(d.get("text_color", "#ffffff"))
    fs = max(6, int(d.get("font_size", 18) * (eh / 450)))

    positions = _diagram_layout(nodes, edges, ew, eh)
    node_w = min(ew / 4.5, 160.0 * (ew / 700))
    node_h = min(eh / 7.0, 50.0 * (eh / 450))

    # Edges
    p.setPen(QPen(line_col, max(1, int(eh * 0.004))))
    p.setBrush(Qt.BrushStyle.NoBrush)
    for a, b in edges:
        if a < len(positions) and b < len(positions) and positions[a] and positions[b]:
            ax, ay = positions[a]; bx, by = positions[b]
            dx = bx - ax; dy = by - ay; dist = math.hypot(dx, dy)
            if dist < 1: continue
            offset = min(node_w, node_h) / 2 + 2
            sx = ax + dx/dist*offset; sy = ay + dy/dist*offset
            ex = bx - dx/dist*offset; ey = by - dy/dist*offset
            p.drawLine(QPointF(sx, sy), QPointF(ex, ey))
            arl = max(8.0, node_h * 0.28); ara = 0.4
            ang = math.atan2(ey-sy, ex-sx)
            p.drawLine(QPointF(ex, ey),
                       QPointF(ex - arl*math.cos(ang-ara), ey - arl*math.sin(ang-ara)))
            p.drawLine(QPointF(ex, ey),
                       QPointF(ex - arl*math.cos(ang+ara), ey - arl*math.sin(ang+ara)))

    # Nodes
    font = QFont("Segoe UI", fs)
    p.setFont(font)
    for i, pos in enumerate(positions):
        if pos is None: continue
        nx, ny = pos
        rect = QRectF(nx - node_w/2, ny - node_h/2, node_w, node_h)
        col = root_col if i == 0 else node_col
        grad = QLinearGradient(nx-node_w/2, ny-node_h/2, nx-node_w/2, ny+node_h/2)
        grad.setColorAt(0, col.lighter(120)); grad.setColorAt(1, col)
        p.setBrush(QBrush(grad))
        p.setPen(QPen(col.lighter(160), max(1, int(node_h*0.04))))
        p.drawRoundedRect(rect, 8, 8)
        p.setPen(text_col)
        lbl = str(nodes[i]) if i < len(nodes) else f"Node {i}"
        p.drawText(rect, Qt.AlignmentFlag.AlignCenter | Qt.TextFlag.TextWordWrap, lbl)


class ShapeElement(BaseElement):
    def paint(self, p, option, widget=None):
        d = self.data
        w, h = float(d["w"]), float(d["h"])
        kind = d["type"]
        p.setOpacity(float(d.get("opacity", 1.0)))
        # Fill / Brush
        if d.get("fill_type") == "gradient":
            brush = _shape_gradient(d, w, h)
        else:
            brush = QBrush(QColor(d.get("fill", "#5294e2")))
        bc = d.get("border_color", "#ffffff")
        bw = int(d.get("border_width", 2))
        p.setBrush(brush)
        p.setPen(QPen(QColor(bc), bw) if bw else Qt.PenStyle.NoPen)

        if kind == ET_RECT:
            r = int(d.get("border_radius", 0))
            if r > 0:
                p.drawRoundedRect(QRectF(0, 0, w, h), r, r)
            else:
                p.drawRect(QRectF(0, 0, w, h))
        elif kind == ET_ELLIPSE:
            p.drawEllipse(QRectF(0, 0, w, h))
        elif kind == ET_LINE:
            p.setPen(QPen(QColor(d.get("color", "#ffffff")),
                          int(d.get("line_width", 3))))
            p.drawLine(QPointF(0, 0), QPointF(w, h))
        elif kind == ET_TRIANGLE:
            poly = QPolygonF([QPointF(w/2, 0),
                              QPointF(w, h),
                              QPointF(0, h)])
            p.drawPolygon(poly)
        elif kind == ET_STAR:
            pts = int(d.get("points", 5))
            path = _build_star_path(w/2, h/2, min(w,h)/2, min(w,h)/4, pts)
            p.drawPath(path)
        elif kind == ET_ARROW:
            path = _build_arrow_path(w, h, d.get("direction", "right"))
            p.drawPath(path)

        p.setOpacity(1.0)
        if self._selected_locally:
            p.setPen(QPen(QColor("#5294e2"), 2, Qt.PenStyle.DashLine))
            p.setBrush(Qt.BrushStyle.NoBrush)
            p.drawRect(QRectF(0, 0, w, h))


class CodeElement(BaseElement):
    def paint(self, p, option, widget=None):
        d = self.data
        w, h = float(d["w"]), float(d["h"])
        bg = QColor(d.get("bg_color", "#1e1e2e"))
        p.fillRect(QRectF(0, 0, w, h), bg)
        # Border
        p.setPen(QPen(QColor("#333344"), 1))
        p.setBrush(Qt.BrushStyle.NoBrush)
        p.drawRect(QRectF(0, 0, w, h))
        # Code text — simple colored rendering
        font = QFont("Consolas", int(d.get("font_size", 18)))
        p.setFont(font)
        fm = p.fontMetrics()
        lh = fm.height() + 2
        import re
        code = d.get("code_text", "")
        y = 8
        kw_color = QColor("#cba6f7")
        str_color = QColor("#a6e3a1")
        cmt_color = QColor("#6c7086")
        num_color = QColor("#fab387")
        default_color = QColor(d.get("text_color", "#cdd6f4"))
        kw_re = re.compile(
            r'\b(def|class|import|from|return|if|elif|else|for|while|try|'
            r'except|finally|with|as|pass|True|False|None|var|let|const|'
            r'function|async|await|new)\b')
        str_re = re.compile(r'(\".*?\"|\'.*?\')')
        cmt_re = re.compile(r'(#[^\n]*|//[^\n]*)')
        num_re = re.compile(r'\b\d+(\.\d+)?\b')
        for line in code.split("\n"):
            if y > h - lh:
                break
            # Determine segments with colors
            segments: list[tuple[int, int, QColor]] = []
            taken = [False] * (len(line) + 1)
            def mark(regex, color):
                for m in regex.finditer(line):
                    s, e = m.start(), m.end()
                    if not any(taken[s:e]):
                        segments.append((s, e, color))
                        for i in range(s, e):
                            taken[i] = True
            mark(cmt_re, cmt_color)
            mark(str_re, str_color)
            mark(kw_re, kw_color)
            mark(num_re, num_color)
            segments.sort(key=lambda x: x[0])
            x = 8
            prev = 0
            for s, e, col in segments:
                if prev < s:
                    p.setPen(default_color)
                    p.drawText(int(x), int(y + fm.ascent()),
                               line[prev:s])
                    x += fm.horizontalAdvance(line[prev:s])
                p.setPen(col)
                p.drawText(int(x), int(y + fm.ascent()), line[s:e])
                x += fm.horizontalAdvance(line[s:e])
                prev = e
            if prev < len(line):
                p.setPen(default_color)
                p.drawText(int(x), int(y + fm.ascent()), line[prev:])
            y += lh
        if self._selected_locally:
            p.setPen(QPen(QColor("#5294e2"), 2, Qt.PenStyle.DashLine))
            p.setBrush(Qt.BrushStyle.NoBrush)
            p.drawRect(QRectF(0, 0, w, h))


class TableElement(BaseElement):
    def paint(self, p, option, widget=None):
        d = self.data
        w, h = float(d["w"]), float(d["h"])
        rows = d.get("rows", 3)
        cols = d.get("cols", 3)
        cells = d.get("cells", [])
        cell_w = w / cols
        cell_h = h / rows
        font = QFont("Segoe UI", int(d.get("font_size", 18)))
        p.setFont(font)
        for r in range(rows):
            for c in range(cols):
                cx = c * cell_w
                cy = r * cell_h
                rect = QRectF(cx, cy, cell_w, cell_h)
                if r == 0:
                    p.fillRect(rect, QColor(d.get("header_bg", "#1a3a5a")))
                    p.setPen(QColor(d.get("header_color", "#ffffff")))
                else:
                    p.fillRect(rect, QColor(d.get("cell_bg", "#1c1c1c")))
                    p.setPen(QColor(d.get("cell_color", "#e0e0e0")))
                # Border
                p.setPen(QPen(QColor(d.get("border_color", "#333333")), 1))
                p.drawRect(rect)
                # Text
                text = ""
                if r < len(cells) and c < len(cells[r]):
                    text = str(cells[r][c])
                if r == 0:
                    p.setPen(QColor(d.get("header_color", "#ffffff")))
                else:
                    p.setPen(QColor(d.get("cell_color", "#e0e0e0")))
                p.drawText(rect.adjusted(4, 2, -4, -2),
                           Qt.AlignmentFlag.AlignCenter |
                           Qt.TextFlag.TextWordWrap, text)
        if self._selected_locally:
            p.setPen(QPen(QColor("#5294e2"), 2, Qt.PenStyle.DashLine))
            p.setBrush(Qt.BrushStyle.NoBrush)
            p.drawRect(QRectF(0, 0, w, h))


class ChartElement(BaseElement):
    def paint(self, p, option, widget=None):
        _draw_chart(p, self.data, float(self.data["w"]), float(self.data["h"]))
        if self._selected_locally:
            p.setPen(QPen(QColor("#5294e2"), 2, Qt.PenStyle.DashLine))
            p.setBrush(Qt.BrushStyle.NoBrush)
            p.drawRect(QRectF(0, 0, self.data["w"], self.data["h"]))

    def mouseDoubleClickEvent(self, e):
        dlg = _ChartDataDialog(self.data, None)
        dlg.setStyleSheet(_STYLE)
        if dlg.exec() == QDialog.DialogCode.Accepted:
            self.update()
            if self.scene_ref:
                self.scene_ref._on_element_changed()


class DiagramElement(BaseElement):
    def paint(self, p, option, widget=None):
        _draw_diagram(p, self.data, float(self.data["w"]), float(self.data["h"]))
        if self._selected_locally:
            p.setPen(QPen(QColor("#5294e2"), 2, Qt.PenStyle.DashLine))
            p.setBrush(Qt.BrushStyle.NoBrush)
            p.drawRect(QRectF(0, 0, self.data["w"], self.data["h"]))

    def mouseDoubleClickEvent(self, e):
        dlg = _DiagramDataDialog(self.data, None)
        dlg.setStyleSheet(_STYLE)
        if dlg.exec() == QDialog.DialogCode.Accepted:
            self.update()
            if self.scene_ref:
                self.scene_ref._on_element_changed()


class _ChartDataDialog(QDialog):
    def __init__(self, data: dict, parent=None):
        super().__init__(parent)
        self.data = data
        self.setWindowTitle("Date grafic")
        self.setMinimumWidth(440)
        layout = QVBoxLayout(self)

        row = QHBoxLayout()
        row.addWidget(QLabel("Tip grafic:"))
        self._type = QComboBox()
        self._type.addItems(["bar", "line", "pie"])
        self._type.setCurrentText(data.get("chart_type", "bar"))
        row.addWidget(self._type, 1)
        layout.addLayout(row)

        row2 = QHBoxLayout()
        row2.addWidget(QLabel("Titlu:"))
        self._title = QLineEdit(data.get("title", ""))
        row2.addWidget(self._title, 1)
        layout.addLayout(row2)

        layout.addWidget(QLabel("Etichete (separate prin virgulă):"))
        self._labels = QLineEdit(", ".join(str(l) for l in data.get("labels", [])))
        layout.addWidget(self._labels)

        layout.addWidget(QLabel("Valori (numere separate prin virgulă):"))
        self._values = QLineEdit(", ".join(
            f"{v:.0f}" if float(v) == int(float(v)) else str(v)
            for v in data.get("values", [])))
        layout.addWidget(self._values)

        btns = QDialogButtonBox(
            QDialogButtonBox.StandardButton.Ok |
            QDialogButtonBox.StandardButton.Cancel)
        btns.accepted.connect(self._accept)
        btns.rejected.connect(self.reject)
        layout.addWidget(btns)

    def _accept(self):
        self.data["chart_type"] = self._type.currentText()
        self.data["title"]  = self._title.text()
        self.data["labels"] = [l.strip() for l in self._labels.text().split(",")
                               if l.strip()]
        vals = []
        for v in self._values.text().split(","):
            try:
                vals.append(float(v.strip()))
            except ValueError:
                pass
        self.data["values"] = vals
        self.accept()


class _DiagramDataDialog(QDialog):
    def __init__(self, data: dict, parent=None):
        super().__init__(parent)
        self.data = data
        self.setWindowTitle("Date diagramă")
        self.setMinimumWidth(440)
        layout = QVBoxLayout(self)

        layout.addWidget(QLabel(
            "Noduri (un nod per linie):"))
        self._nodes = QTextEdit()
        self._nodes.setPlainText("\n".join(str(n) for n in data.get("nodes", [])))
        self._nodes.setMaximumHeight(120)
        layout.addWidget(self._nodes)

        layout.addWidget(QLabel(
            "Conexiuni (format: sursă→destinație, ex: 0→1, una per linie):"))
        self._edges = QTextEdit()
        edges_txt = "\n".join(f"{a}→{b}" for a, b in data.get("edges", []))
        self._edges.setPlainText(edges_txt)
        self._edges.setMaximumHeight(100)
        layout.addWidget(self._edges)

        btns = QDialogButtonBox(
            QDialogButtonBox.StandardButton.Ok |
            QDialogButtonBox.StandardButton.Cancel)
        btns.accepted.connect(self._accept)
        btns.rejected.connect(self.reject)
        layout.addWidget(btns)

    def _accept(self):
        self.data["nodes"] = [l.strip() for l in
                               self._nodes.toPlainText().split("\n") if l.strip()]
        edges = []
        for line in self._edges.toPlainText().split("\n"):
            line = line.strip().replace("->", "→").replace(" ", "")
            if "→" in line:
                parts = line.split("→")
                try:
                    edges.append([int(parts[0]), int(parts[1])])
                except ValueError:
                    pass
        self.data["edges"] = edges
        self.accept()


def _make_element(data: dict, scene_ref) -> BaseElement:
    kind = data.get("type")
    if kind == ET_TEXT:    return TextElement(data, scene_ref)
    if kind == ET_IMAGE:   return ImageElement(data, scene_ref)
    if kind == ET_CODE:    return CodeElement(data, scene_ref)
    if kind == ET_TABLE:   return TableElement(data, scene_ref)
    if kind == ET_CHART:   return ChartElement(data, scene_ref)
    if kind == ET_DIAGRAM: return DiagramElement(data, scene_ref)
    return ShapeElement(data, scene_ref)


# ══════════════════════════════════════════════════════════════════════════════
# Presentation Scene
# ══════════════════════════════════════════════════════════════════════════════

class PresentationScene(QGraphicsScene):
    element_selected  = pyqtSignal(dict)
    element_changed   = pyqtSignal()
    selection_cleared = pyqtSignal()

    def __init__(self, parent=None):
        super().__init__(0, 0, CANVAS_W, CANVAS_H, parent)
        self._slide_data: dict = _default_slide()
        self._items: list[BaseElement] = []
        self._selected_item: BaseElement | None = None
        self._bg_items: list = []
        self._grid_visible  = False
        self._grid_size     = 40
        self._snap_enabled  = False
        self._guide_lines: list = []
        self._draw_background()

    # ── Background ────────────────────────────────────────────────────────────

    def _draw_background(self):
        for it in self._bg_items:
            if it.scene() == self:
                self.removeItem(it)
        self._bg_items.clear()
        d = self._slide_data
        bg_type = d.get("bg_type", "solid")

        if bg_type == "gradient":
            angle = d.get("bg_gradient_angle", 135)
            rad = math.radians(angle)
            cx, cy = CANVAS_W / 2, CANVAS_H / 2
            dx = math.cos(rad) * CANVAS_W / 2
            dy = math.sin(rad) * CANVAS_H / 2
            grad = QLinearGradient(cx - dx, cy - dy, cx + dx, cy + dy)
            grad.setColorAt(0, QColor(d.get("bg_gradient_from", "#1a1a2a")))
            grad.setColorAt(1, QColor(d.get("bg_gradient_to", "#000000")))
            rect = self.addRect(0, 0, CANVAS_W, CANVAS_H,
                                QPen(Qt.PenStyle.NoPen), QBrush(grad))
        else:
            rect = self.addRect(0, 0, CANVAS_W, CANVAS_H,
                                QPen(Qt.PenStyle.NoPen),
                                QBrush(QColor(d.get("bg_color", "#000000"))))
        rect.setZValue(-1000)
        self._bg_items.append(rect)

        bg_image = d.get("bg_image", "")
        if bg_image and os.path.exists(bg_image):
            pix = QPixmap(bg_image).scaled(
                CANVAS_W, CANVAS_H,
                Qt.AspectRatioMode.KeepAspectRatioByExpanding,
                Qt.TransformationMode.SmoothTransformation)
            pix_item = self.addPixmap(pix)
            pix_item.setZValue(-999)
            self._bg_items.append(pix_item)

    def drawBackground(self, painter: QPainter, rect: QRectF):
        super().drawBackground(painter, rect)
        if not self._grid_visible:
            return
        pen = QPen(QColor("#2a2a2a"), 1)
        painter.setPen(pen)
        gs = self._grid_size
        left = int(rect.left()) - (int(rect.left()) % gs)
        top  = int(rect.top())  - (int(rect.top())  % gs)
        x = left
        while x <= rect.right():
            painter.drawLine(x, int(rect.top()), x, int(rect.bottom()))
            x += gs
        y = top
        while y <= rect.bottom():
            painter.drawLine(int(rect.left()), y, int(rect.right()), y)
            y += gs

    def drawForeground(self, painter: QPainter, rect: QRectF):
        if not self._guide_lines:
            return
        pen = QPen(QColor("#ff4444"), 1, Qt.PenStyle.DashLine)
        painter.setPen(pen)
        for line in self._guide_lines:
            painter.drawLine(*line)

    # ── Load / get ────────────────────────────────────────────────────────────

    def load_slide(self, slide_data: dict):
        self.clear()
        self._items.clear()
        self._selected_item = None
        self._bg_items.clear()
        self._slide_data = slide_data
        self._draw_background()
        for el in slide_data.get("elements", []):
            item = _make_element(el, self)
            self.addItem(item)
            self._items.append(item)

    def get_slide_data(self) -> dict:
        return self._slide_data

    # ── Selection ─────────────────────────────────────────────────────────────

    def _deselect_all(self):
        for it in self._items:
            it.set_selected_state(False)
        self._selected_item = None
        self._guide_lines.clear()
        self.update()

    def _select_item(self, item: BaseElement):
        self._deselect_all()
        self._selected_item = item
        item.set_selected_state(True)
        self.element_selected.emit(item.data)

    def mousePressEvent(self, e):
        if e.button() == Qt.MouseButton.LeftButton:
            items_at = [it for it in self.items(e.scenePos())
                        if isinstance(it, BaseElement)]
            if items_at:
                self._select_item(items_at[0])
            else:
                self._deselect_all()
                self.selection_cleared.emit()
        super().mousePressEvent(e)

    # ── Element management ────────────────────────────────────────────────────

    def add_element(self, kind: str) -> BaseElement:
        d = _default_element(kind)
        d["x"] = (CANVAS_W - d["w"]) // 2
        d["y"] = (CANVAS_H - d["h"]) // 2
        d["z"] = len(self._items)
        self._slide_data.setdefault("elements", []).append(d)
        item = _make_element(d, self)
        self.addItem(item)
        self._items.append(item)
        self._select_item(item)
        self.element_changed.emit()
        return item

    def _on_element_changed(self):
        self._update_smart_guides()
        self.element_changed.emit()

    def _on_element_deleted(self, data: dict):
        elements = self._slide_data.get("elements", [])
        if data in elements:
            elements.remove(data)
        self._items = [it for it in self._items if it.data is not data]
        if self._selected_item and self._selected_item.data is data:
            self._selected_item = None
            self.selection_cleared.emit()
        self.element_changed.emit()

    def delete_selected(self):
        if self._selected_item:
            self._selected_item._delete_self()

    def bring_forward(self):
        if self._selected_item:
            self._selected_item._bring_forward()

    def send_back(self):
        if self._selected_item:
            self._selected_item._send_back()

    def _update_smart_guides(self):
        self._guide_lines.clear()
        sel = self._selected_item
        if not sel:
            return
        sx1 = sel.data["x"]; sy1 = sel.data["y"]
        sx2 = sx1 + sel.data["w"]; sy2 = sy1 + sel.data["h"]
        scx = (sx1 + sx2) / 2; scy = (sy1 + sy2) / 2
        SNAP = 8
        for it in self._items:
            if it is sel:
                continue
            ox1 = it.data["x"]; oy1 = it.data["y"]
            ox2 = ox1 + it.data["w"]; oy2 = oy1 + it.data["h"]
            ocx = (ox1 + ox2) / 2; ocy = (oy1 + oy2) / 2
            # Horizontal guides
            for sy, oy in [(sy1, oy1), (sy1, oy2), (sy2, oy1),
                           (sy2, oy2), (scy, ocy)]:
                if abs(sy - oy) < SNAP:
                    self._guide_lines.append(
                        (int(min(sx1,ox1)-20), int(sy),
                         int(max(sx2,ox2)+20), int(sy)))
            # Vertical guides
            for sx, ox in [(sx1, ox1), (sx1, ox2), (sx2, ox1),
                           (sx2, ox2), (scx, ocx)]:
                if abs(sx - ox) < SNAP:
                    self._guide_lines.append(
                        (int(sx), int(min(sy1,oy1)-20),
                         int(sx), int(max(sy2,oy2)+20)))
        self.update()

    # ── Thumbnail ─────────────────────────────────────────────────────────────

    def render_thumbnail(self, w: int, h: int) -> QPixmap:
        pix = QPixmap(w, h)
        pix.fill(QColor(self._slide_data.get("bg_color", "#000000")))
        painter = QPainter(pix)
        self.render(painter, QRectF(0, 0, w, h),
                    QRectF(0, 0, CANVAS_W, CANVAS_H))
        painter.end()
        return pix


# ══════════════════════════════════════════════════════════════════════════════
# Slide View
# ══════════════════════════════════════════════════════════════════════════════

class SlideView(QGraphicsView):
    def __init__(self, scene, parent=None):
        super().__init__(scene, parent)
        self.setRenderHint(QPainter.RenderHint.Antialiasing)
        self.setRenderHint(QPainter.RenderHint.SmoothPixmapTransform)
        self.setDragMode(QGraphicsView.DragMode.NoDrag)
        self.setTransformationAnchor(QGraphicsView.ViewportAnchor.AnchorUnderMouse)
        self.setResizeAnchor(QGraphicsView.ViewportAnchor.AnchorViewCenter)
        self._zoom = 1.0

    def fit_to_view(self):
        self.fitInView(QRectF(0, 0, CANVAS_W, CANVAS_H),
                       Qt.AspectRatioMode.KeepAspectRatio)
        self._zoom = self.transform().m11()

    def wheelEvent(self, e):
        if e.modifiers() & Qt.KeyboardModifier.ControlModifier:
            factor = 1.15 if e.angleDelta().y() > 0 else 1/1.15
            self._zoom = max(0.05, min(self._zoom * factor, 8.0))
            self.setTransform(QTransform().scale(self._zoom, self._zoom))
        else:
            super().wheelEvent(e)

    def resizeEvent(self, e):
        super().resizeEvent(e)
        self.fit_to_view()


# ══════════════════════════════════════════════════════════════════════════════
# Properties Panel
# ══════════════════════════════════════════════════════════════════════════════

class _ColorBtn(QPushButton):
    colorChanged = pyqtSignal(str)

    def __init__(self, color="#ffffff", allow_empty=False, parent=None):
        super().__init__(parent)
        self._color = color
        self._allow_empty = allow_empty
        self.setFixedHeight(26)
        self._refresh()
        self.clicked.connect(self._pick)

    def _refresh(self):
        if self._color:
            self.setStyleSheet(
                f"background:{self._color}; border:1px solid #444; border-radius:3px;")
            self.setText("")
        else:
            self.setStyleSheet("background:#1c1c1c; border:1px solid #333; border-radius:3px;")
            self.setText("(fără)")

    def color(self): return self._color

    def set_color(self, c: str):
        self._color = c
        self._refresh()

    def _pick(self):
        init = QColor(self._color) if self._color else QColor("#ffffff")
        col = QColorDialog.getColor(init, self)
        if col.isValid():
            self._color = col.name()
            self._refresh()
            self.colorChanged.emit(self._color)


class PropertiesPanel(QWidget):
    changed = pyqtSignal(dict)

    def __init__(self, parent=None):
        super().__init__(parent)
        self._data: dict | None = None
        root = QVBoxLayout(self)
        root.setContentsMargins(6, 6, 6, 6)
        root.setSpacing(4)

        self._no_sel_lbl = QLabel("Selectează un element\npentru proprietăți")
        self._no_sel_lbl.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._no_sel_lbl.setStyleSheet("color:#444; font-size:11px;")
        root.addWidget(self._no_sel_lbl)

        self._scroll = QScrollArea()
        self._scroll.setWidgetResizable(True)
        self._scroll.hide()
        root.addWidget(self._scroll, 1)

        self._inner = QWidget()
        self._form = QVBoxLayout(self._inner)
        self._form.setContentsMargins(0, 0, 0, 0)
        self._form.setSpacing(4)
        self._scroll.setWidget(self._inner)

    def clear_selection(self):
        self._data = None
        self._no_sel_lbl.show()
        self._scroll.hide()

    def load_element(self, data: dict):
        self._data = data
        self._no_sel_lbl.hide()
        self._scroll.show()
        self._rebuild(data)

    def _rebuild(self, d: dict):
        # Clear
        while self._form.count():
            item = self._form.takeAt(0)
            if item.widget():
                item.widget().deleteLater()

        kind = d.get("type", "")
        self._section(kind.upper())

        # ── Common: position / size / rotation ──
        self._row("X",       self._spin("x", d, 0, CANVAS_W))
        self._row("Y",       self._spin("y", d, 0, CANVAS_H))
        self._row("Lățime",  self._spin("w", d, 10, CANVAS_W))
        self._row("Înălțime",self._spin("h", d, 10, CANVAS_H))
        self._row("Rotație°",self._spin("rotation", d, 0, 360,
                                        post=self._apply_rotation))
        self._row("Z-index", self._spin("z", d, -100, 100))

        # ── Type-specific ──
        if kind == ET_TEXT:
            self._section("TEXT")
            self._row("Text", self._multiline("text", d))
            self._row("Font", self._font_combo("font", d))
            self._row("Mărime", self._spin("font_size", d, 6, 400))
            self._row("Culoare", self._color("color", d))
            self._row("Bold", self._check("bold", d))
            self._row("Italic", self._check("italic", d))
            self._row("Subliniat", self._check("underline", d))
            self._row("Aliniere", self._combo("align", d,
                                              ["left","center","right"]))
            self._row("Spațiere linii",
                      self._dspin("line_spacing", d, 0.5, 5.0))
            self._row("Spațiere litere",
                      self._spin("letter_spacing", d, 0, 50))
            self._row("Fundal text",
                      self._color("bg_color", d, allow_empty=True))
            self._row("Contur text",
                      self._color("border_color", d, allow_empty=True))

        elif kind == ET_IMAGE:
            self._section("IMAGINE")
            btn = QPushButton("📁 Alege imagine…")
            btn.clicked.connect(self._pick_image)
            self._row("Fișier", btn)
            self._row("Opacitate", self._dspin("opacity", d, 0.0, 1.0))
            self._row("Col. rotunde", self._spin("border_radius", d, 0, 200))

        elif kind in _SHAPE_TYPES:
            self._section("FORMĂ")
            if kind != ET_LINE:
                self._row("Tip umplere",
                          self._combo("fill_type", d,
                                      ["solid","gradient"]))
                self._row("Culoare", self._color("fill", d))
                self._row("Gradient de la",
                          self._color("gradient_from", d))
                self._row("Gradient la",
                          self._color("gradient_to", d))
                self._row("Unghi gradient",
                          self._spin("gradient_angle", d, 0, 360))
            if kind == ET_RECT:
                self._row("Col. rotunde",
                          self._spin("border_radius", d, 0, 200))
            self._row("Culoare contur", self._color("border_color", d))
            self._row("Grosime contur",
                      self._spin("border_width", d, 0, 30))
            self._row("Opacitate",
                      self._dspin("opacity", d, 0.0, 1.0))
            if kind == ET_STAR:
                self._row("Număr vârfuri",
                          self._spin("points", d, 3, 12))
            if kind == ET_ARROW:
                self._row("Direcție",
                          self._combo("direction", d,
                                      ["right","left","up","down"]))
            if kind == ET_LINE:
                self._row("Culoare linie", self._color("color", d))
                self._row("Grosime", self._spin("line_width", d, 1, 30))

        elif kind == ET_CODE:
            self._section("COD SURSĂ")
            self._row("Limbaj",
                      self._combo("language", d, ["python","javascript",
                                                   "html","css","sql"]))
            self._row("Mărime font",
                      self._spin("font_size", d, 8, 60))
            self._row("Fundal", self._color("bg_color", d))

        elif kind == ET_TABLE:
            self._section("TABEL")
            self._row("Rânduri", self._spin("rows", d, 1, 20,
                                            post=self._resize_table))
            self._row("Coloane", self._spin("cols", d, 1, 10,
                                            post=self._resize_table))
            self._row("Fundal header",
                      self._color("header_bg", d))
            self._row("Text header",
                      self._color("header_color", d))
            self._row("Fundal celulă",
                      self._color("cell_bg", d))
            self._row("Bordură", self._color("border_color", d))
            self._row("Mărime font",
                      self._spin("font_size", d, 8, 60))

        elif kind == ET_CHART:
            self._section("GRAFIC")
            edit_btn = QPushButton("📊 Editare date…")
            edit_btn.clicked.connect(self._edit_chart_data)
            self._row("Date", edit_btn)
            self._row("Tip", self._combo("chart_type", d, ["bar", "line", "pie"]))
            self._row("Titlu", self._lineedit("title", d))
            self._row("Culoare bare", self._color("value_color", d))
            self._row("Fundal", self._color("bg_color", d))
            self._row("Text", self._color("text_color", d))
            self._row("Afișează val.", self._check("show_values", d))

        elif kind == ET_DIAGRAM:
            self._section("DIAGRAMĂ")
            edit_btn2 = QPushButton("✏️ Editare noduri…")
            edit_btn2.clicked.connect(self._edit_diagram_data)
            self._row("Noduri", edit_btn2)
            self._row("Col. rădăcină", self._color("root_color", d))
            self._row("Col. noduri",   self._color("node_color", d))
            self._row("Col. linii",    self._color("line_color", d))
            self._row("Text",          self._color("text_color", d))
            self._row("Mărime font",   self._spin("font_size", d, 6, 80))

        # ── Shadow ──
        self._section("UMBRĂ")
        sh = d.setdefault("shadow", {})
        self._row("Activă", self._check_nested(sh, "enabled",
                                               post=self._refresh_shadow))
        self._row("Culoare", self._color_nested(sh, "color"))
        self._row("Blur",  self._spin_nested(sh, "blur", 0, 50))
        self._row("Offset X", self._spin_nested(sh, "offset_x", -30, 30))
        self._row("Offset Y", self._spin_nested(sh, "offset_y", -30, 30))

        # ── Animation ──
        self._section("ANIMAȚIE")
        anim = d.setdefault("animation", {})
        self._row("Intrare",
                  self._combo_nested(anim, "entrance", _ENTRANCES))
        self._row("Ieșire",
                  self._combo_nested(anim, "exit", _EXITS))
        self._row("Loop",
                  self._combo_nested(anim, "loop", _LOOPS))
        self._row("Delay (ms)",
                  self._spin_nested(anim, "delay", 0, 10000))
        self._row("Durată (ms)",
                  self._spin_nested(anim, "duration", 100, 5000))

        self._form.addStretch()

    # ── Widget factories ──────────────────────────────────────────────────────

    def _section(self, title: str):
        lbl = QLabel(f"─── {title}")
        lbl.setStyleSheet(
            "color:#5294e2; font-size:10px; font-weight:700;"
            " letter-spacing:1px; padding-top:6px;")
        self._form.addWidget(lbl)

    def _row(self, label: str, widget: QWidget):
        row = QHBoxLayout()
        row.setSpacing(4)
        lbl = QLabel(label)
        lbl.setStyleSheet("color:#999; font-size:11px;")
        lbl.setFixedWidth(90)
        row.addWidget(lbl)
        row.addWidget(widget, 1)
        self._form.addLayout(row)

    def _spin(self, key, d, mn, mx, post=None):
        w = QSpinBox()
        w.setRange(mn, mx)
        w.setValue(int(d.get(key, 0)))
        def on_change(v):
            d[key] = v
            if post:
                post()
            self.changed.emit(d)
        w.valueChanged.connect(on_change)
        return w

    def _dspin(self, key, d, mn, mx):
        w = QDoubleSpinBox()
        w.setRange(mn, mx)
        w.setSingleStep(0.05)
        w.setDecimals(2)
        w.setValue(float(d.get(key, 1.0)))
        def on_change(v):
            d[key] = round(v, 2)
            self.changed.emit(d)
        w.valueChanged.connect(on_change)
        return w

    def _check(self, key, d):
        w = QCheckBox()
        w.setChecked(bool(d.get(key, False)))
        def on_change(v):
            d[key] = v
            self.changed.emit(d)
        w.toggled.connect(on_change)
        return w

    def _combo(self, key, d, options):
        w = QComboBox()
        for o in options:
            w.addItem(o)
        val = str(d.get(key, options[0]))
        if val in options:
            w.setCurrentText(val)
        def on_change(v):
            d[key] = v
            self.changed.emit(d)
        w.currentTextChanged.connect(on_change)
        return w

    def _color(self, key, d, allow_empty=False):
        btn = _ColorBtn(d.get(key, "#ffffff"), allow_empty)
        def on_change(c):
            d[key] = c
            self.changed.emit(d)
        btn.colorChanged.connect(on_change)
        return btn

    def _multiline(self, key, d):
        w = QTextEdit()
        w.setPlainText(d.get(key, ""))
        w.setMaximumHeight(70)
        def on_change():
            d[key] = w.toPlainText()
            self.changed.emit(d)
        w.textChanged.connect(on_change)
        return w

    def _font_combo(self, key, d):
        w = QFontComboBox()
        w.setCurrentFont(QFont(d.get(key, "Segoe UI")))
        def on_change(f):
            d[key] = f.family()
            self.changed.emit(d)
        w.currentFontChanged.connect(on_change)
        return w

    def _lineedit(self, key, d):
        w = QLineEdit(str(d.get(key, "")))
        def on_change(v):
            d[key] = v
            self.changed.emit(d)
        w.textChanged.connect(on_change)
        return w

    def _spin_nested(self, container, key, mn, mx):
        w = QSpinBox()
        w.setRange(mn, mx)
        w.setValue(int(container.get(key, 0)))
        def on_change(v):
            container[key] = v
            if self._data:
                self.changed.emit(self._data)
        w.valueChanged.connect(on_change)
        return w

    def _check_nested(self, container, key, post=None):
        w = QCheckBox()
        w.setChecked(bool(container.get(key, False)))
        def on_change(v):
            container[key] = v
            if post:
                post()
            if self._data:
                self.changed.emit(self._data)
        w.toggled.connect(on_change)
        return w

    def _color_nested(self, container, key):
        btn = _ColorBtn(container.get(key, "#000000"))
        def on_change(c):
            container[key] = c
            if self._data:
                self.changed.emit(self._data)
        btn.colorChanged.connect(on_change)
        return btn

    def _combo_nested(self, container, key, options):
        w = QComboBox()
        for o in options:
            w.addItem(o)
        val = str(container.get(key, options[0]))
        if val in options:
            w.setCurrentText(val)
        def on_change(v):
            container[key] = v
            if self._data:
                self.changed.emit(self._data)
        w.currentTextChanged.connect(on_change)
        return w

    def _apply_rotation(self):
        """Called when rotation spinbox changes — needs scene refresh."""
        if self._data:
            self.changed.emit(self._data)

    def _refresh_shadow(self):
        if self._data:
            self.changed.emit(self._data)

    def _resize_table(self):
        if not self._data:
            return
        rows = self._data.get("rows", 3)
        cols = self._data.get("cols", 3)
        cells = self._data.get("cells", [])
        new_cells = []
        for r in range(rows):
            row_data = cells[r] if r < len(cells) else []
            new_row = []
            for c in range(cols):
                new_row.append(row_data[c] if c < len(row_data) else "")
            new_cells.append(new_row)
        self._data["cells"] = new_cells
        self.changed.emit(self._data)

    def _pick_image(self):
        path, _ = QFileDialog.getOpenFileName(
            self, "Alege imagine", "",
            "Imagini (*.png *.jpg *.jpeg *.bmp *.gif *.webp)")
        if path and self._data:
            self._data["path"] = path
            self.changed.emit(self._data)

    def _edit_chart_data(self):
        if not self._data:
            return
        dlg = _ChartDataDialog(self._data, self)
        dlg.setStyleSheet(_STYLE)
        if dlg.exec() == QDialog.DialogCode.Accepted:
            self.changed.emit(self._data)

    def _edit_diagram_data(self):
        if not self._data:
            return
        dlg = _DiagramDataDialog(self._data, self)
        dlg.setStyleSheet(_STYLE)
        if dlg.exec() == QDialog.DialogCode.Accepted:
            self.changed.emit(self._data)


# ══════════════════════════════════════════════════════════════════════════════
# Animation Timeline
# ══════════════════════════════════════════════════════════════════════════════

class AnimationTimeline(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setMaximumHeight(140)
        layout = QVBoxLayout(self)
        layout.setContentsMargins(6, 4, 6, 4)
        layout.setSpacing(2)

        hdr = QLabel("TIMELINE ANIMAȚII")
        hdr.setStyleSheet(
            "color:#5294e2; font-size:10px; font-weight:700; letter-spacing:2px;")
        layout.addWidget(hdr)

        self._container = QWidget()
        self._vlay = QVBoxLayout(self._container)
        self._vlay.setContentsMargins(0, 0, 0, 0)
        self._vlay.setSpacing(2)
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setWidget(self._container)
        scroll.setMaximumHeight(100)
        layout.addWidget(scroll, 1)

    def load_slide(self, slide_data: dict):
        while self._vlay.count():
            item = self._vlay.takeAt(0)
            if item.widget():
                item.widget().deleteLater()

        animated = [el for el in slide_data.get("elements", [])
                    if el.get("animation", {}).get("entrance", "none") != "none"]

        if not animated:
            lbl = QLabel("Nicio animație pe acest slide")
            lbl.setStyleSheet("color:#444; font-size:11px; padding:4px;")
            self._vlay.addWidget(lbl)
        else:
            for el in animated:
                anim = el.get("animation", {})
                row = QWidget()
                rl = QHBoxLayout(row)
                rl.setContentsMargins(2, 2, 2, 2)
                rl.setSpacing(6)
                kind = el.get("type", "?")
                name = (el.get("text","")[:18] + "…"
                        if kind == ET_TEXT and el.get("text")
                        else kind)
                n = QLabel(name)
                n.setFixedWidth(130)
                n.setStyleSheet("color:#ccc; font-size:11px;")
                rl.addWidget(n)
                e_lbl = QLabel(anim.get("entrance","none"))
                e_lbl.setFixedWidth(100)
                e_lbl.setStyleSheet("color:#5294e2; font-size:11px;")
                rl.addWidget(e_lbl)
                delay = int(anim.get("delay", 0))
                dur   = int(anim.get("duration", 500))
                t = QLabel(f"delay {delay}ms · dur {dur}ms")
                t.setStyleSheet("color:#666; font-size:10px;")
                rl.addWidget(t)
                rl.addStretch()
                self._vlay.addWidget(row)

        self._vlay.addStretch()


# ══════════════════════════════════════════════════════════════════════════════
# Slide List Panel
# ══════════════════════════════════════════════════════════════════════════════

class SlideListPanel(QWidget):
    slide_selected   = pyqtSignal(int)
    slide_added      = pyqtSignal()
    slide_deleted    = pyqtSignal(int)
    slide_duplicated = pyqtSignal(int)

    THUMB_W, THUMB_H = 168, 95

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setFixedWidth(200)
        layout = QVBoxLayout(self)
        layout.setContentsMargins(4, 4, 4, 4)
        layout.setSpacing(4)

        hdr = QLabel("SLIDE-URI")
        hdr.setStyleSheet(
            "color:#5294e2; font-size:10px; font-weight:700; letter-spacing:2px;")
        layout.addWidget(hdr)

        self._list = QListWidget()
        self._list.setDragDropMode(QListWidget.DragDropMode.InternalMove)
        self._list.setSpacing(2)
        self._list.currentRowChanged.connect(self.slide_selected)
        self._list.setContextMenuPolicy(
            Qt.ContextMenuPolicy.CustomContextMenu)
        self._list.customContextMenuRequested.connect(self._ctx)
        layout.addWidget(self._list, 1)

        add_btn = QPushButton("➕ Slide nou")
        add_btn.clicked.connect(self.slide_added)
        layout.addWidget(add_btn)

    def populate(self, slides: list[dict], current: int = 0):
        self._list.blockSignals(True)
        self._list.clear()
        for i, s in enumerate(slides):
            item = QListWidgetItem(f"  Slide {i+1}")
            item.setBackground(QColor(s.get("bg_color", "#000000")))
            item.setForeground(QColor("#e0e0e0"))
            self._list.addItem(item)
        self._list.blockSignals(False)
        if 0 <= current < self._list.count():
            self._list.setCurrentRow(current)

    def update_thumbnail(self, idx: int, pix: QPixmap):
        item = self._list.item(idx)
        if item:
            item.setIcon(pix.scaled(
                self.THUMB_W, self.THUMB_H,
                Qt.AspectRatioMode.KeepAspectRatio,
                Qt.TransformationMode.SmoothTransformation))
            self._list.setIconSize(QSize(self.THUMB_W, self.THUMB_H))

    def _ctx(self, pos):
        from PyQt6.QtWidgets import QMenu
        idx = self._list.currentRow()
        menu = QMenu(self)
        menu.addAction("🗑 Șterge",  lambda: self.slide_deleted.emit(idx))
        menu.addAction("📋 Duplică", lambda: self.slide_duplicated.emit(idx))
        menu.exec(self._list.mapToGlobal(pos))


# ══════════════════════════════════════════════════════════════════════════════
# Slide Notes Panel
# ══════════════════════════════════════════════════════════════════════════════

class SlideNotesPanel(QWidget):
    notes_changed = pyqtSignal(str)

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setMaximumHeight(100)
        layout = QVBoxLayout(self)
        layout.setContentsMargins(4, 2, 4, 2)
        layout.setSpacing(2)
        hdr = QLabel("NOTE PREZENTATOR")
        hdr.setStyleSheet(
            "color:#5294e2; font-size:10px; font-weight:700; letter-spacing:2px;")
        layout.addWidget(hdr)
        self._edit = QTextEdit()
        self._edit.setPlaceholderText("Note vizibile doar pe Stage…")
        self._edit.setMaximumHeight(70)
        self._edit.textChanged.connect(
            lambda: self.notes_changed.emit(self._edit.toPlainText()))
        layout.addWidget(self._edit)

    def load(self, notes: str):
        self._edit.blockSignals(True)
        self._edit.setPlainText(notes or "")
        self._edit.blockSignals(False)


# ══════════════════════════════════════════════════════════════════════════════
# Background Dialog
# ══════════════════════════════════════════════════════════════════════════════

class BackgroundDialog(QDialog):
    def __init__(self, slide_data: dict, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Fundal Slide")
        self.setFixedSize(440, 320)
        self.setStyleSheet(_STYLE)
        self.sd = slide_data

        layout = QVBoxLayout(self)
        form = QFormLayout()
        form.setSpacing(8)

        # Type
        self._type_combo = QComboBox()
        for t in ["solid","gradient","image"]:
            self._type_combo.addItem(t)
        self._type_combo.setCurrentText(slide_data.get("bg_type","solid"))
        self._type_combo.currentTextChanged.connect(self._on_type)
        form.addRow("Tip:", self._type_combo)

        # Solid color
        self._color_btn = _ColorBtn(slide_data.get("bg_color","#000000"))
        form.addRow("Culoare solidă:", self._color_btn)

        # Gradient
        self._grad_from = _ColorBtn(
            slide_data.get("bg_gradient_from","#1a1a2a"))
        self._grad_to   = _ColorBtn(
            slide_data.get("bg_gradient_to","#000000"))
        self._grad_angle = QSpinBox()
        self._grad_angle.setRange(0, 360)
        self._grad_angle.setValue(int(
            slide_data.get("bg_gradient_angle", 135)))
        form.addRow("Gradient de la:", self._grad_from)
        form.addRow("Gradient la:",    self._grad_to)
        form.addRow("Unghi gradient:", self._grad_angle)

        # Image
        self._img_edit = QLineEdit(slide_data.get("bg_image",""))
        img_btn = QPushButton("📁")
        img_btn.setFixedWidth(32)
        img_btn.clicked.connect(self._pick_image)
        row = QHBoxLayout()
        row.addWidget(self._img_edit)
        row.addWidget(img_btn)
        form.addRow("Imagine fundal:", row)

        # Transition
        self._trans_combo = QComboBox()
        for t in _TRANSITIONS:
            self._trans_combo.addItem(t)
        self._trans_combo.setCurrentText(
            slide_data.get("transition","fade"))
        form.addRow("Tranziție:", self._trans_combo)

        self._trans_ms = QSpinBox()
        self._trans_ms.setRange(0, 3000)
        self._trans_ms.setSingleStep(50)
        self._trans_ms.setValue(int(slide_data.get("transition_ms",400)))
        form.addRow("Durată tranziție (ms):", self._trans_ms)

        layout.addLayout(form)
        layout.addStretch()

        bb = QDialogButtonBox(
            QDialogButtonBox.StandardButton.Ok |
            QDialogButtonBox.StandardButton.Cancel)
        bb.accepted.connect(self._apply)
        bb.rejected.connect(self.reject)
        layout.addWidget(bb)
        self._on_type(self._type_combo.currentText())

    def _on_type(self, t: str):
        is_grad  = t == "gradient"
        is_solid = t == "solid"
        is_image = t == "image"

    def _pick_image(self):
        path, _ = QFileDialog.getOpenFileName(
            self, "Alege imagine fundal", "",
            "Imagini (*.png *.jpg *.jpeg *.bmp *.webp)")
        if path:
            self._img_edit.setText(path)

    def _pick_color(self):
        col = QColorDialog.getColor(
            QColor(self.sd.get("bg_color","#000000")), self)
        if col.isValid():
            self.sd["bg_color"] = col.name()
            self._color_btn.set_color(col.name())

    def _apply(self):
        self.sd["bg_type"]           = self._type_combo.currentText()
        self.sd["bg_color"]          = self._color_btn.color()
        self.sd["bg_gradient_from"]  = self._grad_from.color()
        self.sd["bg_gradient_to"]    = self._grad_to.color()
        self.sd["bg_gradient_angle"] = self._grad_angle.value()
        self.sd["bg_image"]          = self._img_edit.text().strip()
        self.sd["transition"]        = self._trans_combo.currentText()
        self.sd["transition_ms"]     = self._trans_ms.value()
        self.accept()


# ══════════════════════════════════════════════════════════════════════════════
# Templates Dialog
# ══════════════════════════════════════════════════════════════════════════════

def _build_template(key: str) -> list[dict]:
    def slide(bg="#000000", **kw):
        s = _default_slide()
        s["bg_color"] = bg
        s.update(kw)
        return s

    def txt(text, x, y, w, h, size=48, bold=False, color="#ffffff",
            align="center"):
        e = _default_element(ET_TEXT)
        e.update({"text": text, "x": x, "y": y, "w": w, "h": h,
                  "font_size": size, "bold": bold, "color": color,
                  "align": align})
        return e

    if key == "title":
        s = slide(bg_type="gradient", bg_gradient_from="#1a2a4a",
                  bg_gradient_to="#000000", bg_gradient_angle=135)
        s["elements"] = [
            txt("Titlu Prezentare", 160, 330, 1600, 160, 72, True),
            txt("Subtitlu sau descriere", 300, 520, 1320, 80,
                36, color="#aaaacc"),
        ]
        return [s]

    elif key == "title_content":
        s = slide("#0a0a14")
        s["elements"] = [
            txt("Titlu Slide", 60, 40, 1800, 120, 54, True),
            txt("Conținut principal\n• Punct 1\n• Punct 2\n• Punct 3",
                60, 200, 1800, 700, 32, color="#dddddd", align="left"),
        ]
        return [s]

    elif key == "two_column":
        s = slide("#0a0a14")
        s["elements"] = [
            txt("Titlu", 60, 40, 1800, 100, 48, True),
            txt("Coloana stângă\n• Element A\n• Element B",
                60, 180, 860, 700, 28, color="#dddddd", align="left"),
            txt("Coloana dreaptă\n• Element X\n• Element Y",
                1000, 180, 860, 700, 28, color="#dddddd", align="left"),
        ]
        return [s]

    elif key == "quote":
        s = slide(bg_type="gradient", bg_gradient_from="#1a0a2a",
                  bg_gradient_to="#000000", bg_gradient_angle=135)
        s["elements"] = [
            txt('"Citat inspirațional lung care ocupă\nmajoritatea slide-ului"',
                120, 280, 1680, 380, 52, color="#e9c46a"),
            txt("— Autor, Sursă", 120, 700, 1680, 60, 28,
                color="#888888"),
        ]
        return [s]

    elif key == "worship":
        s = slide(bg_type="gradient", bg_gradient_from="#000000",
                  bg_gradient_to="#0a0a1a", bg_gradient_angle=180)
        s["elements"] = [
            txt("Versul 1 — Linia 1\nLinia 2\nLinia 3\nLinia 4",
                100, 250, 1720, 580, 56, color="#ffffff",
                align="center"),
        ]
        return [s]

    else:  # blank
        return [_default_slide()]


class TemplatesDialog(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Șabloane prezentare")
        self.setFixedSize(580, 360)
        self.setStyleSheet(_STYLE)
        self._selected_key: str | None = None

        layout = QVBoxLayout(self)
        layout.addWidget(QLabel(
            "Alege un șablon de pornire:"))

        grid = QGridLayout()
        grid.setSpacing(10)

        templates = [
            ("blank",         "Gol",           "#111111"),
            ("title",         "Titlu",          "#1a2a4a"),
            ("title_content", "Titlu + Conținut","#0a0a14"),
            ("two_column",    "Două Coloane",    "#0a0a14"),
            ("quote",         "Citat",           "#1a0a2a"),
            ("worship",       "Worship",         "#000000"),
        ]
        self._btns = {}
        for i, (key, name, color) in enumerate(templates):
            pix = QPixmap(160, 90)
            pix.fill(QColor(color))
            p = QPainter(pix)
            p.setPen(QColor("#ffffff"))
            p.drawText(QRect(0, 0, 160, 90),
                       Qt.AlignmentFlag.AlignCenter, name)
            p.end()
            btn = QPushButton()
            btn.setFixedSize(160, 90)
            btn.setIcon(pix)
            btn.setIconSize(QSize(156, 86))
            btn.setCheckable(True)
            btn.clicked.connect(lambda c, k=key, b=btn: self._select(k, b))
            self._btns[key] = btn
            lbl = QLabel(name)
            lbl.setAlignment(Qt.AlignmentFlag.AlignCenter)
            cell = QWidget()
            cl = QVBoxLayout(cell)
            cl.setSpacing(3)
            cl.addWidget(btn)
            cl.addWidget(lbl)
            grid.addWidget(cell, i // 3, i % 3)

        layout.addLayout(grid)

        bb = QDialogButtonBox(
            QDialogButtonBox.StandardButton.Ok |
            QDialogButtonBox.StandardButton.Cancel)
        bb.accepted.connect(self.accept)
        bb.rejected.connect(self.reject)
        layout.addWidget(bb)

        self._select("blank", self._btns["blank"])

    def _select(self, key: str, btn: QPushButton):
        for b in self._btns.values():
            b.setChecked(False)
            b.setStyleSheet("")
        btn.setChecked(True)
        btn.setStyleSheet("border:3px solid #5294e2; border-radius:6px;")
        self._selected_key = key

    def get_slides(self) -> list[dict]:
        return _build_template(self._selected_key or "blank")


# ══════════════════════════════════════════════════════════════════════════════
# PPTX Importer
# ══════════════════════════════════════════════════════════════════════════════

class PPTXImporter:
    EMU_W = 9144000
    EMU_H = 5143500

    @classmethod
    def import_file(cls, path: str) -> list[dict]:
        try:
            from pptx import Presentation as PptxPres
        except ImportError:
            raise ImportError(
                "python-pptx nu este instalat.\nRulați: pip install python-pptx")

        prs = PptxPres(path)
        sw = prs.slide_width  or cls.EMU_W
        sh = prs.slide_height or cls.EMU_H
        sx = CANVAS_W / sw
        sy = CANVAS_H / sh
        slides = []

        for pptx_slide in prs.slides:
            slide = _default_slide()

            # Background color (best-effort)
            try:
                bg = pptx_slide.background.fill
                if bg.type is not None:
                    from pptx.util import Pt
                    col = bg.fore_color.rgb
                    slide["bg_color"] = f"#{col}"
            except Exception:
                pass

            # Elements
            for shape in pptx_slide.shapes:
                try:
                    el = None
                    ex = int(shape.left   * sx)
                    ey = int(shape.top    * sy)
                    ew = int(shape.width  * sx)
                    eh = int(shape.height * sy)

                    if shape.has_text_frame:
                        el = _default_element(ET_TEXT)
                        el.update({"x": ex, "y": ey, "w": ew, "h": eh})
                        texts = [p.text for p in shape.text_frame.paragraphs
                                 if p.text.strip()]
                        el["text"] = "\n".join(texts)
                        # Font from first run
                        try:
                            para = shape.text_frame.paragraphs[0]
                            if para.runs:
                                run = para.runs[0]
                                if run.font.size:
                                    el["font_size"] = max(
                                        8, int(run.font.size.pt * sy * 2))
                                el["bold"]   = bool(run.font.bold)
                                el["italic"] = bool(run.font.italic)
                                try:
                                    el["color"] = f"#{run.font.color.rgb}"
                                except Exception:
                                    pass
                        except Exception:
                            pass

                    elif hasattr(shape, "image"):
                        import tempfile, os as _os
                        img = shape.image
                        ext = img.ext
                        tmp = tempfile.NamedTemporaryFile(
                            delete=False, suffix=f".{ext}")
                        tmp.write(img.blob)
                        tmp.close()
                        el = _default_element(ET_IMAGE)
                        el.update({"x": ex, "y": ey, "w": ew, "h": eh,
                                   "path": tmp.name})

                    if el:
                        el["z"] = len(slide["elements"])
                        slide["elements"].append(el)
                except Exception:
                    pass

            slides.append(slide)

        return slides if slides else [_default_slide()]


# ══════════════════════════════════════════════════════════════════════════════
# Main Presentation Editor Window
# ══════════════════════════════════════════════════════════════════════════════

class PresentationEditor(QMainWindow):
    saved = pyqtSignal(int, str, list)

    def __init__(self, pres_id=None, title="", slides=None, parent=None):
        super().__init__(parent)
        self.setMinimumSize(1280, 720)
        self.setStyleSheet(_STYLE)

        self._pres_id = pres_id
        if slides:
            raw = slides
            self._slides = (json.loads(raw)
                            if isinstance(raw, str) else list(raw))
        else:
            self._slides = [_default_slide()]

        self._presentation_title = title or "Prezentare nouă"
        self.setWindowTitle(
            f"Cantio — Editor Prezentări — {self._presentation_title}")

        self._current_slide = 0
        self._undo_stack: list[list] = []
        self._redo_stack: list[list] = []
        self._modified   = False

        self._scene = PresentationScene()
        self._scene.element_selected.connect(self._on_element_selected)
        self._scene.element_changed.connect(self._on_element_changed)
        self._scene.selection_cleared.connect(self._on_selection_cleared)

        self._build_ui()
        self._load_presentation()

    # ── UI ────────────────────────────────────────────────────────────────────

    def _build_ui(self):
        tb = QToolBar("Elemente", self)
        tb.setMovable(False)
        self.addToolBar(tb)
        self._build_toolbar(tb)

        central = QWidget()
        self.setCentralWidget(central)
        root = QHBoxLayout(central)
        root.setContentsMargins(0, 0, 0, 0)
        root.setSpacing(0)

        # Left: slide list
        self._slide_panel = SlideListPanel()
        self._slide_panel.slide_selected.connect(self._switch_slide)
        self._slide_panel.slide_added.connect(self._add_slide)
        self._slide_panel.slide_deleted.connect(self._delete_slide)
        self._slide_panel.slide_duplicated.connect(self._duplicate_slide)
        root.addWidget(self._slide_panel)

        # Centre: canvas + notes + timeline
        centre = QWidget()
        centre_lay = QVBoxLayout(centre)
        centre_lay.setContentsMargins(0, 0, 0, 0)
        centre_lay.setSpacing(0)

        self._view = SlideView(self._scene)
        centre_lay.addWidget(self._view, 1)

        self._notes = SlideNotesPanel()
        self._notes.notes_changed.connect(self._on_notes_changed)
        centre_lay.addWidget(self._notes)

        self._timeline = AnimationTimeline()
        centre_lay.addWidget(self._timeline)

        root.addWidget(centre, 1)

        # Right: properties
        self._props = PropertiesPanel()
        self._props.setFixedWidth(270)
        self._props.changed.connect(self._on_props_changed)
        root.addWidget(self._props)

        self._update_slide_panel()

    def _build_toolbar(self, tb: QToolBar):
        S = ("QPushButton{background:#1c1c1c;color:#ccc;"
             "border:1px solid #2a2a2a;border-radius:4px;padding:5px 9px;}"
             "QPushButton:hover{background:#252525;color:#fff;}"
             "QPushButton:checked{background:#1a3a5c;"
             "color:#5294e2;border-color:#5294e2;}")

        def btn(label, slot, tip=""):
            b = QPushButton(label)
            b.setToolTip(tip)
            b.setStyleSheet(S)
            b.clicked.connect(slot)
            tb.addWidget(b)
            return b

        btn("💾 Salvează",   self._save,              "Ctrl+S")
        btn("📤 Export",     self._export_json,        "Export .json")
        btn("📥 Import JSON",self._import_json,        "Import .json")
        btn("📊 Import PPTX",self._import_pptx,       "Import PowerPoint")
        tb.addSeparator()

        btn("📋 Șabloane",  self._open_templates,     "Șabloane slide")
        tb.addSeparator()

        btn("T Text",       lambda: self._add_el(ET_TEXT),    "Adaugă text")
        btn("□ Drept",      lambda: self._add_el(ET_RECT),    "Dreptunghi")
        btn("○ Elipsă",     lambda: self._add_el(ET_ELLIPSE), "Elipsă")
        btn("△ Triunghi",   lambda: self._add_el(ET_TRIANGLE),"Triunghi")
        btn("★ Stea",       lambda: self._add_el(ET_STAR),    "Stea")
        btn("➡ Săgeată",    lambda: self._add_el(ET_ARROW),   "Săgeată")
        btn("— Linie",      lambda: self._add_el(ET_LINE),    "Linie")
        btn("🖼 Imagine",   self._add_image,                  "Imagine")
        btn("</> Cod",      lambda: self._add_el(ET_CODE),    "Bloc cod")
        btn("⊞ Tabel",      lambda: self._add_el(ET_TABLE),   "Tabel")
        btn("📊 Grafic",    lambda: self._add_el(ET_CHART),   "Grafic (bar/line/pie)")
        btn("🔗 Diagramă",  lambda: self._add_el(ET_DIAGRAM), "Diagramă/Organigram")
        tb.addSeparator()

        btn("⬆ Față",    self._scene.bring_forward, "Aduce în față")
        btn("⬇ Spate",   self._scene.send_back,     "Trimite în spate")
        btn("🗑 Șterge",  self._scene.delete_selected,"Delete")
        tb.addSeparator()

        btn("🖼 Fundal",  self._edit_background, "Fundal slide")

        # Grid toggle
        self._grid_btn = QPushButton("# Grid")
        self._grid_btn.setCheckable(True)
        self._grid_btn.setStyleSheet(S)
        self._grid_btn.toggled.connect(self._toggle_grid)
        tb.addWidget(self._grid_btn)

        # Snap toggle
        self._snap_btn = QPushButton("⊞ Snap")
        self._snap_btn.setCheckable(True)
        self._snap_btn.setStyleSheet(S)
        self._snap_btn.toggled.connect(self._toggle_snap)
        tb.addWidget(self._snap_btn)

        tb.addSeparator()
        btn("↩ Undo", self._undo, "Ctrl+Z")
        btn("↪ Redo", self._redo, "Ctrl+Y")

        spacer = QWidget()
        spacer.setSizePolicy(QSizePolicy.Policy.Expanding,
                             QSizePolicy.Policy.Preferred)
        tb.addWidget(spacer)
        self._zoom_lbl = QLabel("100%")
        self._zoom_lbl.setStyleSheet(
            "color:#555; font-size:11px; padding-right:10px;")
        tb.addWidget(self._zoom_lbl)

        # Shortcuts
        from PyQt6.QtGui import QShortcut
        QShortcut(QKeySequence("Ctrl+S"), self).activated.connect(self._save)
        QShortcut(QKeySequence("Ctrl+Z"), self).activated.connect(self._undo)
        QShortcut(QKeySequence("Ctrl+Y"), self).activated.connect(self._redo)
        QShortcut(QKeySequence("Delete"), self).activated.connect(
            self._scene.delete_selected)
        QShortcut(QKeySequence("Ctrl+D"), self).activated.connect(
            self._duplicate_selected)

    # ── Grid / Snap ───────────────────────────────────────────────────────────

    def _toggle_grid(self, on: bool):
        self._scene._grid_visible = on
        self._scene.update()

    def _toggle_snap(self, on: bool):
        self._scene._snap_enabled = on

    # ── Slides ────────────────────────────────────────────────────────────────

    def _switch_slide(self, idx: int):
        if idx < 0 or idx >= len(self._slides):
            return
        self._push_snapshot()
        self._current_slide = idx
        self._scene.load_slide(self._slides[idx])
        self._notes.load(self._slides[idx].get("notes", ""))
        self._timeline.load_slide(self._slides[idx])
        self._props.clear_selection()

    def _add_slide(self):
        self._push_snapshot()
        new_s = _default_slide()
        self._slides.insert(self._current_slide + 1, new_s)
        self._current_slide += 1
        self._update_slide_panel()
        self._scene.load_slide(new_s)
        self._notes.load("")
        self._timeline.load_slide(new_s)

    def _delete_slide(self, idx: int):
        if len(self._slides) <= 1:
            QMessageBox.information(self, "Info",
                                    "Trebuie să existe cel puțin un slide.")
            return
        self._push_snapshot()
        self._slides.pop(idx)
        self._current_slide = max(
            0, min(self._current_slide, len(self._slides) - 1))
        self._update_slide_panel()
        self._scene.load_slide(self._slides[self._current_slide])

    def _duplicate_slide(self, idx: int):
        if 0 <= idx < len(self._slides):
            self._push_snapshot()
            dup = copy.deepcopy(self._slides[idx])
            self._slides.insert(idx + 1, dup)
            self._update_slide_panel()

    def _update_slide_panel(self):
        self._slide_panel.populate(self._slides, self._current_slide)
        QTimer.singleShot(60, self._refresh_thumbnail)

    def _refresh_thumbnail(self):
        try:
            pix = self._scene.render_thumbnail(168, 95)
            self._slide_panel.update_thumbnail(self._current_slide, pix)
        except Exception:
            pass

    # ── Elements ──────────────────────────────────────────────────────────────

    def _add_el(self, kind: str):
        self._push_snapshot()
        self._scene.add_element(kind)

    def _add_image(self):
        path, _ = QFileDialog.getOpenFileName(
            self, "Alege imagine", "",
            "Imagini (*.png *.jpg *.jpeg *.bmp *.gif *.webp)")
        if path:
            self._push_snapshot()
            item = self._scene.add_element(ET_IMAGE)
            item.data["path"] = path
            item._load_pixmap()
            item.update()

    def _duplicate_selected(self):
        sel = self._scene._selected_item
        if sel:
            sel._duplicate_self()

    def _on_element_selected(self, data: dict):
        self._props.load_element(data)

    def _on_element_changed(self):
        self._modified = True
        self._refresh_thumbnail()
        self._timeline.load_slide(self._slides[self._current_slide])

    def _on_selection_cleared(self):
        self._props.clear_selection()

    def _on_props_changed(self, data: dict):
        self._modified = True
        for item in self._scene._items:
            if item.data is data:
                item.prepareGeometryChange()
                item._update_handle_positions()
                item.setPos(data.get("x", 0), data.get("y", 0))
                rot = data.get("rotation", 0)
                item.setTransformOriginPoint(
                    data["w"] / 2, data["h"] / 2)
                item.setRotation(rot)
                op = data.get("opacity", 1.0)
                if data.get("type") not in (ET_RECT, ET_ELLIPSE,
                                            ET_TRIANGLE, ET_STAR,
                                            ET_ARROW, ET_LINE):
                    item.setOpacity(float(op))
                item._apply_shadow()
                if isinstance(item, ImageElement):
                    item._load_pixmap()
                item.update()
                break
        self._refresh_thumbnail()

    def _on_notes_changed(self, text: str):
        self._slides[self._current_slide]["notes"] = text
        self._modified = True

    # ── Background ────────────────────────────────────────────────────────────

    def _edit_background(self):
        dlg = BackgroundDialog(self._slides[self._current_slide], self)
        if dlg.exec() == QDialog.DialogCode.Accepted:
            self._scene.load_slide(self._slides[self._current_slide])
            self._notes.load(
                self._slides[self._current_slide].get("notes", ""))
            self._modified = True
            self._refresh_thumbnail()

    # ── Templates ─────────────────────────────────────────────────────────────

    def _open_templates(self):
        dlg = TemplatesDialog(self)
        if dlg.exec() == QDialog.DialogCode.Accepted:
            slides = dlg.get_slides()
            if slides:
                self._push_snapshot()
                self._slides = slides
                self._current_slide = 0
                self._update_slide_panel()
                self._scene.load_slide(self._slides[0])
                self._notes.load(self._slides[0].get("notes", ""))
                self._timeline.load_slide(self._slides[0])

    # ── PPTX Import ───────────────────────────────────────────────────────────

    def _import_pptx(self):
        path, _ = QFileDialog.getOpenFileName(
            self, "Import PowerPoint", "",
            "PowerPoint (*.pptx *.ppt)")
        if not path:
            return
        try:
            slides = PPTXImporter.import_file(path)
            if slides:
                self._push_snapshot()
                self._slides = slides
                self._current_slide = 0
                self._update_slide_panel()
                self._scene.load_slide(self._slides[0])
                self._notes.load("")
                self._timeline.load_slide(self._slides[0])
                QMessageBox.information(
                    self, "Import PPTX",
                    f"Importate {len(slides)} slide-uri din:\n{path}")
        except ImportError as e:
            QMessageBox.warning(self, "Import PPTX", str(e))
        except Exception as e:
            QMessageBox.critical(self, "Eroare import", str(e))

    # ── Undo / Redo ───────────────────────────────────────────────────────────

    def _push_snapshot(self):
        snap = copy.deepcopy(self._slides)
        self._undo_stack.append(snap)
        if len(self._undo_stack) > MAX_UNDO:
            self._undo_stack.pop(0)
        self._redo_stack.clear()

    def _undo(self):
        if not self._undo_stack:
            return
        self._redo_stack.append(copy.deepcopy(self._slides))
        self._slides = self._undo_stack.pop()
        self._current_slide = min(self._current_slide,
                                  len(self._slides) - 1)
        self._update_slide_panel()
        self._scene.load_slide(self._slides[self._current_slide])

    def _redo(self):
        if not self._redo_stack:
            return
        self._undo_stack.append(copy.deepcopy(self._slides))
        self._slides = self._redo_stack.pop()
        self._current_slide = min(self._current_slide,
                                  len(self._slides) - 1)
        self._update_slide_panel()
        self._scene.load_slide(self._slides[self._current_slide])

    # ── Save / Load ───────────────────────────────────────────────────────────

    def _load_presentation(self):
        if self._pres_id is None:
            self._scene.load_slide(self._slides[0])
            self._notes.load(self._slides[0].get("notes", ""))
            self._timeline.load_slide(self._slides[0])
            return
        try:
            import database as db
            pres = db.get_presentation(self._pres_id)
            if pres:
                self.setWindowTitle(
                    f"Cantio — {pres.get('title','Prezentare')}")
                raw = pres.get("slides", pres.get("slides_json", []))
                slides = (json.loads(raw)
                          if isinstance(raw, str) else raw)
                if slides:
                    self._slides = slides
                    self._current_slide = 0
                    self._update_slide_panel()
                    self._scene.load_slide(self._slides[0])
                    self._notes.load(self._slides[0].get("notes",""))
                    self._timeline.load_slide(self._slides[0])
        except Exception as e:
            print(f"[PRES EDITOR] Load error: {e}")

    def _save(self):
        try:
            import database as db
            if self._pres_id is None:
                title, ok = QInputDialog.getText(
                    self, "Salvează", "Titlu prezentare:",
                    text=self._presentation_title)
                if not ok or not title.strip():
                    return
                title = title.strip()
                self._pres_id = db.add_presentation(title, self._slides)
                self._presentation_title = title
                self.setWindowTitle(
                    f"Cantio — Editor Prezentări — {title}")
            else:
                pres = db.get_presentation(self._pres_id)
                title = (pres.get("title", self._presentation_title)
                         if pres else self._presentation_title)
                db.update_presentation(self._pres_id, title, self._slides)
                self._presentation_title = title
            self._modified = False
            self.saved.emit(
                self._pres_id, self._presentation_title, self._slides)
            try:
                from toast_notifications import show_toast
                show_toast("✅ Prezentare salvată", "success")
            except Exception:
                pass
        except Exception as e:
            QMessageBox.critical(self, "Eroare", f"Salvare eșuată:\n{e}")

    def _export_json(self):
        path, _ = QFileDialog.getSaveFileName(
            self, "Export prezentare", "prezentare.json", "JSON (*.json)")
        if path:
            try:
                with open(path, "w", encoding="utf-8") as f:
                    json.dump({"slides": self._slides}, f,
                              ensure_ascii=False, indent=2)
                QMessageBox.information(self, "Export",
                                        f"Salvat în:\n{path}")
            except Exception as e:
                QMessageBox.critical(self, "Eroare", str(e))

    def _import_json(self):
        path, _ = QFileDialog.getOpenFileName(
            self, "Import prezentare", "", "JSON (*.json)")
        if path:
            try:
                with open(path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                slides = data.get("slides",
                                  data if isinstance(data, list) else [])
                if slides:
                    self._push_snapshot()
                    self._slides = slides
                    self._current_slide = 0
                    self._update_slide_panel()
                    self._scene.load_slide(self._slides[0])
                    self._notes.load(self._slides[0].get("notes",""))
                    self._timeline.load_slide(self._slides[0])
            except Exception as e:
                QMessageBox.critical(self, "Eroare import", str(e))

    def closeEvent(self, event):
        if self._modified:
            r = QMessageBox.question(
                self, "Modificări nesalvate",
                "Prezentarea are modificări nesalvate. Salvați?",
                QMessageBox.StandardButton.Save |
                QMessageBox.StandardButton.Discard |
                QMessageBox.StandardButton.Cancel)
            if r == QMessageBox.StandardButton.Save:
                self._save()
            elif r == QMessageBox.StandardButton.Cancel:
                event.ignore()
                return
        event.accept()


# ── Backward-compat alias ─────────────────────────────────────────────────────
PresentationEditorWindow = PresentationEditor


# ══════════════════════════════════════════════════════════════════════════════
# Static renderer (thumbnails + live display via control_window)
# ══════════════════════════════════════════════════════════════════════════════

def render_slide_to_pixmap(slide_data: dict, w: int, h: int) -> QPixmap:
    pix = QPixmap(w, h)
    pix.fill(QColor(slide_data.get("bg_color", "#000000")))
    p = QPainter(pix)
    p.setRenderHint(QPainter.RenderHint.Antialiasing)
    p.setRenderHint(QPainter.RenderHint.TextAntialiasing)

    # Background
    bg_type = slide_data.get("bg_type", "solid")
    if bg_type == "gradient":
        angle = slide_data.get("bg_gradient_angle", 135)
        rad = math.radians(angle)
        cx, cy = w / 2, h / 2
        dx, dy = math.cos(rad) * w / 2, math.sin(rad) * h / 2
        grad = QLinearGradient(cx-dx, cy-dy, cx+dx, cy+dy)
        grad.setColorAt(0, QColor(slide_data.get("bg_gradient_from","#1a1a2a")))
        grad.setColorAt(1, QColor(slide_data.get("bg_gradient_to","#000000")))
        p.fillRect(0, 0, w, h, QBrush(grad))
    elif bg_type == "image":
        bg_image = slide_data.get("bg_image", "")
        if bg_image and os.path.exists(bg_image):
            bg_pix = QPixmap(bg_image).scaled(
                w, h, Qt.AspectRatioMode.KeepAspectRatioByExpanding,
                Qt.TransformationMode.SmoothTransformation)
            p.drawPixmap(0, 0, bg_pix)

    sx = w / CANVAS_W
    sy = h / CANVAS_H

    for el in sorted(slide_data.get("elements", []),
                     key=lambda e: e.get("z", 0)):
        if not el.get("visible", True):
            continue
        kind = el.get("type")
        ex = el.get("x", 0) * sx
        ey = el.get("y", 0) * sy
        ew = el.get("w", 100) * sx
        eh = el.get("h", 40)  * sy
        rot = el.get("rotation", 0)

        p.save()
        if rot:
            p.translate(ex + ew/2, ey + eh/2)
            p.rotate(rot)
            p.translate(-(ew/2), -(eh/2))
            ex_r, ey_r = 0.0, 0.0
        else:
            p.translate(ex, ey)
            ex_r, ey_r = 0.0, 0.0

        p.setOpacity(float(el.get("opacity", 1.0)))

        if kind == ET_TEXT:
            font = QFont(el.get("font","Segoe UI"),
                         max(1, int(el.get("font_size",48) * sx)))
            font.setBold(el.get("bold", False))
            font.setItalic(el.get("italic", False))
            p.setFont(font)
            p.setPen(QColor(el.get("color","#ffffff")))
            align_map = {"left": Qt.AlignmentFlag.AlignLeft,
                         "center": Qt.AlignmentFlag.AlignHCenter,
                         "right":  Qt.AlignmentFlag.AlignRight}
            align = align_map.get(el.get("align","center"),
                                   Qt.AlignmentFlag.AlignHCenter)
            flags = align | Qt.AlignmentFlag.AlignVCenter | Qt.TextFlag.TextWordWrap
            p.drawText(QRectF(ex_r, ey_r, ew, eh), flags,
                       el.get("text",""))

        elif kind == ET_IMAGE:
            path = el.get("path","")
            if path and os.path.exists(path):
                img_pix = QPixmap(path).scaled(
                    int(ew), int(eh),
                    Qt.AspectRatioMode.KeepAspectRatio,
                    Qt.TransformationMode.SmoothTransformation)
                p.drawPixmap(int(ex_r), int(ey_r), img_pix)

        elif kind in _SHAPE_TYPES:
            if el.get("fill_type") == "gradient":
                brush = _shape_gradient(el, ew, eh)
            else:
                brush = QBrush(QColor(el.get("fill","#5294e2")))
            bc = el.get("border_color","#ffffff")
            bw = max(0, int(el.get("border_width",2) * sx))
            p.setBrush(brush)
            p.setPen(QPen(QColor(bc), bw) if bw else Qt.PenStyle.NoPen)
            rect = QRectF(ex_r, ey_r, ew, eh)
            if kind == ET_RECT:
                r = int(el.get("border_radius",0) * sx)
                p.drawRoundedRect(rect, r, r) if r else p.drawRect(rect)
            elif kind == ET_ELLIPSE:
                p.drawEllipse(rect)
            elif kind == ET_TRIANGLE:
                poly = QPolygonF([QPointF(ew/2, 0),
                                  QPointF(ew, eh), QPointF(0, eh)])
                p.drawPolygon(poly)
            elif kind == ET_STAR:
                pts = int(el.get("points",5))
                path = _build_star_path(ew/2, eh/2,
                                        min(ew,eh)/2, min(ew,eh)/4, pts)
                p.drawPath(path)
            elif kind == ET_ARROW:
                path = _build_arrow_path(ew, eh, el.get("direction","right"))
                p.drawPath(path)
            elif kind == ET_LINE:
                lw = max(1, int(el.get("line_width",3) * sx))
                p.setPen(QPen(QColor(el.get("color","#ffffff")), lw))
                p.drawLine(QPointF(ex_r, ey_r),
                           QPointF(ex_r + ew, ey_r + eh))

        elif kind == ET_CODE:
            p.fillRect(QRectF(ex_r, ey_r, ew, eh),
                       QColor(el.get("bg_color","#1e1e2e")))
            font = QFont("Consolas", max(1, int(el.get("font_size",18)*sx)))
            p.setFont(font)
            p.setPen(QColor(el.get("text_color","#cdd6f4")))
            p.drawText(QRectF(ex_r+4, ey_r+4, ew-8, eh-8),
                       Qt.AlignmentFlag.AlignTop | Qt.AlignmentFlag.AlignLeft |
                       Qt.TextFlag.TextWordWrap,
                       el.get("code_text",""))

        elif kind == ET_TABLE:
            rows = el.get("rows",3)
            cols = el.get("cols",3)
            cells = el.get("cells",[])
            cw = ew / cols
            ch = eh / rows
            font = QFont("Segoe UI", max(1,int(el.get("font_size",18)*sx)))
            p.setFont(font)
            for r in range(rows):
                for c in range(cols):
                    rx = ex_r + c * cw
                    ry = ey_r + r * ch
                    cell_rect = QRectF(rx, ry, cw, ch)
                    p.fillRect(cell_rect, QColor(
                        el.get("header_bg","#1a3a5a") if r==0
                        else el.get("cell_bg","#1c1c1c")))
                    p.setPen(QPen(QColor(el.get("border_color","#333")),1))
                    p.drawRect(cell_rect)
                    text = ""
                    if r < len(cells) and c < len(cells[r]):
                        text = str(cells[r][c])
                    p.setPen(QColor(
                        el.get("header_color","#fff") if r==0
                        else el.get("cell_color","#e0e0e0")))
                    p.drawText(cell_rect.adjusted(2,1,-2,-1),
                               Qt.AlignmentFlag.AlignCenter, text)

        elif kind == ET_CHART:
            _draw_chart(p, el, ew, eh)

        elif kind == ET_DIAGRAM:
            _draw_diagram(p, el, ew, eh)

        p.setOpacity(1.0)
        p.restore()

    p.end()
    return pix
